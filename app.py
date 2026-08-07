import base64
import glob
import io
import json
import os
import queue
import re
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import threading
import time
import urllib.error
import urllib.request
import uuid
from datetime import datetime, timedelta
from zoneinfo import ZoneInfo


def _load_env_file():
    env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
    if not os.path.exists(env_path):
        return
    with open(env_path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, _, value = line.partition("=")
            os.environ.setdefault(key.strip(), value.strip().strip('"').strip("'"))


_load_env_file()

from flask import (
    Flask,
    flash,
    jsonify,
    redirect,
    render_template,
    request,
    send_file,
    url_for,
)
from flask_login import (
    LoginManager,
    UserMixin,
    current_user,
    login_required,
    login_user,
    logout_user,
)
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy import text as sa_text
from werkzeug.security import check_password_hash, generate_password_hash

from lessons_data import LESSONS, LONG_QUIZZES
from changelog_data import CHANGELOG
from labs_data import LABS
SECTION_ORDER = ["Fundamentals", "Advanced Fundamentals", "Advanced Database System", "SQL Injection"]
SECTION_LESSON_MINUTES = {"Fundamentals": 40, "Advanced Database System": 45}
AVAILABLE_LESSONS = [l for l in LESSONS if not l.get("coming_soon")]


def lesson_estimate_minutes(lesson):
    base = lesson.get(
        "minutes",
        SECTION_LESSON_MINUTES.get(lesson.get("section", "Fundamentals"), 45),
    )
    videos = lesson.get("videos") or []
    video_minutes = 0
    if isinstance(videos, dict):
        for items in videos.values():
            for v in items:
                video_minutes += v.get("duration", 0)
    else:
        for v in videos:
            video_minutes += v.get("duration", 0)
    return base + video_minutes

LAB_POINTS_PER_LAB = 2
LAB_LESSON_IDS = set(LABS.keys())
CUSTOM_LAB_ROUTES = {19: "sql_lab", 20: "sql_lab_2"}
COINS_PER_LESSON = 10
DIFFICULTY_UNLOCK_COST = {
    "easy": 10,
    "medium": 20,
    "hard": 30,
}
LESSON_DIFFICULTY = {
    1: "easy",
    29: "easy",
    2: "easy",
    3: "easy",
    4: "medium",
    5: "medium",
    6: "medium",
    7: "medium",
    8: "medium",
    9: "medium",
    10: "hard",
    11: "medium",
    12: "medium",
    13: "hard",
    14: "medium",
    15: "medium",
    16: "medium",
    17: "hard",
    18: "hard",
    19: "hard",
    20: "hard",
    21: "hard",
    22: "hard",
    23: "hard",
    24: "hard",
    25: "medium",
    26: "medium",
    27: "hard",
    28: "medium",
    30: "medium",
    31: "medium",
    32: "medium",
    33: "medium",
    34: "medium",
}


def lesson_difficulty(lesson_id):
    return LESSON_DIFFICULTY.get(lesson_id, "medium")


def lesson_unlock_cost(lesson_id):
    return DIFFICULTY_UNLOCK_COST.get(lesson_difficulty(lesson_id), 20)

app = Flask(__name__)
app.config["SECRET_KEY"] = os.environ.get("SECRET_KEY", "change-this-secret-key-in-production")

LOCAL_TZ = ZoneInfo("Asia/Manila")


@app.template_filter("localtime")
def localtime_filter(dt, fmt="%Y-%m-%d %H:%M:%S"):
    if dt is None:
        return ""
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=ZoneInfo("UTC"))
    return dt.astimezone(LOCAL_TZ).strftime(fmt)


APP_VERSION = "1.62"
EREN_VERSION = "2.2"


@app.context_processor
def inject_avatar():
    def avatar_src(user):
        if not user or not user.avatar:
            return ""
        if user.avatar.startswith("data:"):
            return user.avatar
        return url_for("static", filename="uploads/" + user.avatar)

    return {"avatar_src": avatar_src, "app_version": APP_VERSION, "eren_version": EREN_VERSION, "latest_update": CHANGELOG[0]}


@app.context_processor
def inject_coin_settings():
    return {"coins_per_lesson": COINS_PER_LESSON}

DATABASE_URL = os.environ.get("DATABASE_URL", "sqlite:///students.db")
if DATABASE_URL.startswith("postgres://"):
    DATABASE_URL = DATABASE_URL.replace("postgres://", "postgresql://", 1)
app.config["SQLALCHEMY_DATABASE_URI"] = DATABASE_URL
if DATABASE_URL.startswith("postgresql"):
    connect_args = {"sslmode": "require"}
    app.config["SQLALCHEMY_ENGINE_OPTIONS"] = {
        "pool_pre_ping": True,
        "pool_recycle": 240,
        "connect_args": connect_args,
    }

ALLOWED_AVATAR_EXT = {"png", "jpg", "jpeg", "gif", "webp"}
AVATAR_DIR = os.path.join(app.root_path, "static", "uploads")
AVATAR_MIME = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
}

db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = "login"
login_manager.login_message = "Please log in to access that page."
login_manager.login_message_category = "warning"


class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(50), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(200), nullable=False)
    avatar = db.Column(db.Text, nullable=True)
    lab_completed = db.Column(db.Integer, default=0)
    is_admin = db.Column(db.Boolean, default=False, nullable=False)
    is_banned = db.Column(db.Boolean, default=False, nullable=False)
    ban_reason = db.Column(db.String(200), nullable=True)
    coins = db.Column(db.Integer, default=0, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    last_seen = db.Column(db.DateTime, default=datetime.utcnow)
    progress = db.relationship("LessonProgress", backref="user", cascade="all, delete-orphan")
    quiz_scores = db.relationship("QuizScore", backref="user", cascade="all, delete-orphan")
    section_quiz_scores = db.relationship("SectionQuizScore", backref="user", cascade="all, delete-orphan")
    unlocked_lessons = db.relationship("UnlockedLesson", backref="user", cascade="all, delete-orphan")


class LessonProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    lesson_id = db.Column(db.Integer, nullable=False)
    completed_at = db.Column(db.DateTime, default=datetime.utcnow)


class UnlockedLesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    lesson_id = db.Column(db.Integer, nullable=False)
    unlocked_at = db.Column(db.DateTime, default=datetime.utcnow)


class QuizScore(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    lesson_id = db.Column(db.Integer, nullable=False)
    best_score = db.Column(db.Integer, nullable=False)
    total_questions = db.Column(db.Integer, nullable=False)
    attempts = db.Column(db.Integer, default=1, nullable=False)
    last_updated = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)


class SectionQuizScore(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    section = db.Column(db.String(100), nullable=False)
    best_score = db.Column(db.Integer, nullable=False)
    total_questions = db.Column(db.Integer, nullable=False)
    attempts = db.Column(db.Integer, default=1, nullable=False)
    last_answers = db.Column(db.Text, nullable=True)
    topic_ratings = db.Column(db.Text, nullable=True)
    last_updated = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)


class LabProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    lesson_id = db.Column(db.Integer, nullable=False)
    completed_at = db.Column(db.DateTime, default=datetime.utcnow)


class ActivityLog(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    action = db.Column(db.String(100), nullable=False)
    details = db.Column(db.String(500), nullable=True)
    ip_address = db.Column(db.String(45), nullable=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow, index=True)

    user = db.relationship("User", backref=db.backref("activity_logs", lazy="dynamic"))


class FailedLogin(db.Model):
    __tablename__ = "failed_login"
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(120), nullable=True)
    ip_address = db.Column(db.String(45), nullable=True, index=True)
    user_agent = db.Column(db.String(250), nullable=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow, index=True)


class Feedback(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=True)
    name = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(120), nullable=False)
    category = db.Column(db.String(50), nullable=False, default="general")
    message = db.Column(db.Text, nullable=False)
    admin_reply = db.Column(db.Text, nullable=True)
    is_read = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow, index=True)
    replied_at = db.Column(db.DateTime, nullable=True)

    user = db.relationship("User", backref=db.backref("feedbacks", lazy="dynamic"))

    @property
    def like_count(self):
        return self.likes.count()


class FeedbackLike(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    feedback_id = db.Column(db.Integer, db.ForeignKey("feedback.id"), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

    __table_args__ = (db.UniqueConstraint("feedback_id", "user_id", name="uq_feedback_like"),)

    feedback = db.relationship(
        "Feedback", backref=db.backref("likes", lazy="dynamic", cascade="all, delete-orphan")
    )
    user = db.relationship("User", backref=db.backref("feedback_likes", lazy="dynamic"))


def log_activity(user, action, details=""):
    try:
        db.session.add(
            ActivityLog(
                user_id=user.id,
                action=action,
                details=details[:500] if details else None,
                ip_address=request.headers.get("X-Forwarded-For", request.remote_addr or "").split(",")[0].strip() or None,
            )
        )
        db.session.commit()
    except Exception:
        db.session.rollback()


LOGIN_LOCK_ATTEMPTS = 5
LOGIN_LOCK_MINUTES = 30


def _client_ip():
    return (
        request.headers.get("X-Forwarded-For", request.remote_addr or "")
        .split(",")[0]
        .strip()
        or None
    )


def failed_login_count(ip, minutes=LOGIN_LOCK_MINUTES):
    if not ip:
        return 0
    since = datetime.utcnow() - timedelta(minutes=minutes)
    return FailedLogin.query.filter(
        FailedLogin.ip_address == ip, FailedLogin.created_at >= since
    ).count()


def is_ip_locked(ip):
    return failed_login_count(ip) >= LOGIN_LOCK_ATTEMPTS


def ip_lock_remaining_minutes(ip):
    if not is_ip_locked(ip):
        return 0
    since = datetime.utcnow() - timedelta(minutes=LOGIN_LOCK_MINUTES)
    recent = (
        FailedLogin.query.filter(
            FailedLogin.ip_address == ip, FailedLogin.created_at >= since
        )
        .order_by(FailedLogin.created_at.asc())
        .all()
    )
    if len(recent) < LOGIN_LOCK_ATTEMPTS:
        return 0
    fifth = recent[LOGIN_LOCK_ATTEMPTS - 1]
    expires = fifth.created_at + timedelta(minutes=LOGIN_LOCK_MINUTES)
    seconds = max(0, int((expires - datetime.utcnow()).total_seconds()))
    return (seconds + 59) // 60 or 1


BADWORDS = {
    "fuck", "fucking", "fucked", "fuckin", "fck", "motherfucker", "mf",
    "shit", "shitt", "bullshit", "bitch", "bitchy", "bastard", "asshole",
    "a$$hole", "arsehole", "dick", "dickhead", "cunt", "pussy", "twat",
    "wanker", "bollocks", "retard", "retarded", "nigger", "nigga", "faggot",
    "fag", "dyke", "whore", "slut", "hoe", "prostitute", "porn", "porno",
    "pornography", "sex", "sexy", "xxx", "penis", "vagina", "dildo",
    "crap", "damn", "damnit", "hell", "piss", "pissed", "douche", "douchebag",
    "jackass", "jerk", "idiot", "moron", "stupid", "dumbass", "dumb",
    "bobo", "bobo mo", "tangina", "tang ina", "putangina", "putang ina",
    "puta", "puta ka", "ulol", "gago", "gaga", "bwisit", "bwisit ka",
    "tarantado", "bobo", "tanga", "inutil", "walang kwenta", "hayop",
    "hayup", "kupal", "sira ulo", "bobo ka", "animal ka", "peste",
    "unggoy", "baboy", "tamad", "bobo inutil", "epal", "mapangmata",
    "buwisit", "sira", "lintik", "leche", "peste ka", "gagong", "ulol ka",
    "pakingshet", "pakyu", "puking ina", "pokpok", "bading", "bakla",
    "bading ka", "bayot", "puti ng", "inutil ka", "walanghiya",
    "walang hiya", "sira ulo ka", "kainin mo", "mamatay ka", "mamatay",
    "tarantado ka", "demonyo", "demonyo ka", "impakto", "halimaw",
    "palalo", "manyak", "manyak", "manyakis", "pervert", "malibog",
    "libog", "kantot", "kantutin", "kantutan", "tite", "titi", "pekpek",
    "puke", "betlog", "itlog mo", "hinayupak", "anak ng puta", "pota",
    "punyeta", "punyemas", "bilat", "puday", "poyet", "yawa", "yawa ka",
    "pisting yawa", "giatay", "gago ka", "bogo", "hanggaw", "piste",
    "burikat", "landi", "malandi", "makati", "pasaway", "sutil", "astig ka",
}

BADWORD_PATTERN = re.compile(
    r"(?<![a-zA-Z0-9])("
    + "|".join(re.escape(w) for w in sorted(BADWORDS, key=len, reverse=True))
    + r")(?![a-zA-Z0-9])",
    re.IGNORECASE,
)

FEEDBACK_SPAM_LIMIT = 3
FEEDBACK_SPAM_MINUTES = 5


def contains_badword(text):
    return bool(BADWORD_PATTERN.search(text or ""))


def detected_badword(text):
    m = BADWORD_PATTERN.search(text or "")
    return m.group(1) if m else None


def ban_user(user, reason):
    if user is None or user.is_banned or user.is_admin:
        return False
    user.is_banned = True
    user.ban_reason = (reason or "Violation of community rules")[:200]
    db.session.commit()
    log_activity(user, "Account banned", "Reason: %s" % user.ban_reason)
    return True


def unban_user(user):
    if user is None or not user.is_banned:
        return False
    user.is_banned = False
    user.ban_reason = None
    db.session.commit()
    log_activity(user, "Account unbanned", "Restored by admin")
    return True


def record_failed_login(username):
    try:
        db.session.add(
            FailedLogin(
                username=(username or "")[:120],
                ip_address=_client_ip(),
                user_agent=(request.user_agent.string[:250] if request.user_agent else None),
            )
        )
        db.session.commit()
    except Exception:
        db.session.rollback()


def clear_failed_logins(ip=None):
    try:
        if ip:
            FailedLogin.query.filter_by(ip_address=ip).delete(synchronize_session=False)
        else:
            FailedLogin.query.delete(synchronize_session=False)
        db.session.commit()
    except Exception:
        db.session.rollback()


def admin_required(view):
    from functools import wraps

    @wraps(view)
    def wrapped(*args, **kwargs):
        if not current_user.is_authenticated:
            return redirect(url_for("login", next=request.path))
        if not current_user.is_admin:
            flash("You do not have permission to view that page.", "danger")
            return redirect(url_for("dashboard"))
        return view(*args, **kwargs)

    return wrapped


@login_manager.user_loader
def load_user(user_id):
    try:
        return db.session.get(User, int(user_id))
    except Exception:
        return None


MAINTENANCE_MODE = os.environ.get("MAINTENANCE_MODE", "0") == "1"


@app.before_request
def check_maintenance_mode():
    if MAINTENANCE_MODE:
        maintenance_paths = {"/maintenance", "/static/"}
        path = request.path
        if path.startswith("/static/") or path == "/maintenance":
            return None
        return render_template("maintenance.html"), 503


@app.after_request
def add_maintenance_header(response):
    if MAINTENANCE_MODE:
        if request.path != "/maintenance" and not request.path.startswith("/static/"):
            response.headers["Retry-After"] = "3600"
    return response


@app.before_request
def update_last_seen():
    try:
        if current_user.is_authenticated:
            current_user.last_seen = datetime.utcnow()
            db.session.commit()
    except Exception:
        db.session.rollback()


def is_lab_completed(user, lesson_id):
    if user is None:
        return False
    return (
        db.session.query(LabProgress)
        .filter_by(user_id=user.id, lesson_id=lesson_id)
        .first()
        is not None
    )


def get_completed_lab_ids(user):
    if user is None:
        return set()
    return {lp.lesson_id for lp in LabProgress.query.filter_by(user_id=user.id).all()}


def mark_lab_completed(user, lesson_id):
    if not is_lab_completed(user, lesson_id):
        db.session.add(LabProgress(user_id=user.id, lesson_id=lesson_id))
        db.session.commit()


def unmark_lab_completed(user, lesson_id):
    row = (
        db.session.query(LabProgress)
        .filter_by(user_id=user.id, lesson_id=lesson_id)
        .first()
    )
    if row:
        db.session.delete(row)
        db.session.commit()


def get_completed_ids(user):
    if user is None:
        return set()
    return {p.lesson_id for p in user.progress}


def get_unlocked_ids(user):
    if user is None:
        return set()
    return {ul.lesson_id for ul in user.unlocked_lessons}


def get_quiz_scores(user):
    if user is None:
        return {}
    return {qs.lesson_id: qs for qs in user.quiz_scores}


def get_leaderboard():
    rows = db.session.query(QuizScore, User.username).join(
        User, QuizScore.user_id == User.id
    ).all()
    stats = {}
    for qs, username in rows:
        entry = stats.setdefault(
            qs.user_id,
            {
                "username": username,
                "points": 0,
                "quizzes": 0,
                "attempts": 0,
                "labs": 0,
                "lab_points": 0,
            },
        )
        entry["points"] += qs.best_score
        entry["quizzes"] += 1
        entry["attempts"] += qs.attempts
    lab_rows = (
        db.session.query(LabProgress, User.username)
        .join(User, LabProgress.user_id == User.id)
        .all()
    )
    for lp, username in lab_rows:
        entry = stats.setdefault(
            lp.user_id,
            {
                "username": username,
                "points": 0,
                "quizzes": 0,
                "attempts": 0,
                "labs": 0,
                "lab_points": 0,
            },
        )
        entry["labs"] += 1
        entry["lab_points"] += LAB_POINTS_PER_LAB
    for entry in stats.values():
        entry["points"] += entry["lab_points"]
    board = sorted(stats.values(), key=lambda e: (-e["points"], e["attempts"]))
    for i, entry in enumerate(board, start=1):
        entry["rank"] = i
    return board


def get_lesson_labs_leaderboard(lesson_id):
    rows = (
        db.session.query(User.username, db.func.count(LabProgress.id))
        .outerjoin(
            LabProgress,
            db.and_(LabProgress.user_id == User.id, LabProgress.lesson_id == lesson_id),
        )
        .group_by(User.id, User.username)
        .all()
    )
    board = sorted(
        (
            {
                "username": username,
                "labs": labs_done,
                "points": labs_done * LAB_POINTS_PER_LAB,
            }
            for username, labs_done in rows
            if labs_done > 0
        ),
        key=lambda e: (-e["labs"], e["username"].lower()),
    )
    for i, entry in enumerate(board, start=1):
        entry["rank"] = i
    return board


def get_lesson_leaderboard(lesson_id):
    rows = db.session.query(QuizScore, User.username).join(
        User, QuizScore.user_id == User.id
    ).filter(QuizScore.lesson_id == lesson_id).all()
    board = sorted(
        (
            {
                "username": username,
                "best_score": qs.best_score,
                "total": qs.total_questions,
            }
            for qs, username in rows
        ),
        key=lambda e: (-e["best_score"], e["username"].lower()),
    )
    for i, entry in enumerate(board, start=1):
        entry["rank"] = i
    return board


def password_errors(password):
    errors = []
    if len(password) < 8:
        errors.append("at least 8 characters")
    if not re.search(r"[A-Z]", password):
        errors.append("one uppercase letter")
    if not re.search(r"[a-z]", password):
        errors.append("one lowercase letter")
    if not re.search(r"\d", password):
        errors.append("one number")
    if not re.search(r"[^A-Za-z0-9]", password):
        errors.append("one special character")
    return errors


@app.context_processor
def inject_user_stats():
    if current_user.is_authenticated:
        return {"user_done": len(get_completed_ids(current_user)), "user_total": len(AVAILABLE_LESSONS)}
    return {}


@app.after_request
def no_cache_html(response):
    if response.content_type and response.content_type.startswith("text/html"):
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
    return response


@app.before_request
def block_banned_users():
    if current_user.is_authenticated and getattr(current_user, "is_banned", False):
        logout_user()
        flash(
            "Your account has been banned (%s). Contact an admin to get unbanned."
            % (getattr(current_user, "ban_reason", None) or "violation of community rules"),
            "danger",
        )
        return redirect(url_for("login"))


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/changelog")
def changelog():
    return render_template("changelog.html", changelog=CHANGELOG)


@app.route("/faq")
def faq():
    return render_template("faq.html")


@app.route("/feedback", methods=["GET", "POST"])
def feedback():
    if request.method == "POST":
        if not current_user.is_authenticated:
            flash("Please log in to submit feedback.", "warning")
            return redirect(url_for("login", next=request.path))

        name = request.form.get("name", "").strip()
        email = request.form.get("email", "").strip()
        category = request.form.get("category", "general").strip()
        message = request.form.get("message", "").strip()

        if not name or not email or not message:
            flash("Please fill in all required fields.", "warning")
            return redirect(url_for("feedback"))

        bad = detected_badword(message)
        if bad:
            ban_user(
                current_user,
                "Used inappropriate language (\"%s\") in feedback" % bad,
            )
            log_activity(current_user, "Banned by system", "Bad word detected in feedback: \"%s\"" % bad)
            logout_user()
            flash(
                "Your feedback was removed because it contained inappropriate language. "
                "Your account has been banned. Contact an admin to get unbanned.",
                "danger",
            )
            return redirect(url_for("login"))

        since = datetime.utcnow() - timedelta(minutes=FEEDBACK_SPAM_MINUTES)
        recent_count = Feedback.query.filter(
            Feedback.user_id == current_user.id,
            Feedback.created_at >= since,
        ).count()
        if recent_count >= FEEDBACK_SPAM_LIMIT:
            ban_user(
                current_user,
                "Feedback spam (%d submissions in %d minutes)" % (
                    recent_count + 1, FEEDBACK_SPAM_MINUTES,
                ),
            )
            log_activity(current_user, "Banned by system", "Spam: %d feedbacks in %d minutes" % (recent_count + 1, FEEDBACK_SPAM_MINUTES))
            logout_user()
            flash(
                "You are submitting feedback too quickly. This looks like spam, "
                "so your account has been banned. Contact an admin to get unbanned.",
                "danger",
            )
            return redirect(url_for("login"))

        fb = Feedback(
            user_id=current_user.id,
            name=name,
            email=email,
            category=category,
            message=message,
        )
        db.session.add(fb)
        db.session.commit()
        log_activity(current_user, "Submitted feedback", f"Category: {category}")
        flash("Thank you for your feedback! We appreciate it.", "success")
        return redirect(url_for("feedback"))

    feedbacks = Feedback.query.order_by(Feedback.created_at.desc()).limit(50).all()
    liked_ids = set()
    if current_user.is_authenticated:
        liked_ids = {
            l.feedback_id for l in FeedbackLike.query.filter_by(user_id=current_user.id).all()
        }
    return render_template("feedback.html", feedbacks=feedbacks, liked_ids=liked_ids)


@app.route("/feedback/<int:fid>/like", methods=["POST"])
@login_required
def feedback_like(fid):
    fb = db.session.get(Feedback, fid)
    if not fb:
        return jsonify({"error": "Feedback not found."}), 404
    existing = FeedbackLike.query.filter_by(feedback_id=fid, user_id=current_user.id).first()
    if existing:
        db.session.delete(existing)
        db.session.commit()
        return jsonify({"liked": False, "count": fb.like_count})
    db.session.add(FeedbackLike(feedback_id=fid, user_id=current_user.id))
    db.session.commit()
    log_activity(current_user, "Liked feedback", f"Feedback #{fid} from {fb.name}")
    return jsonify({"liked": True, "count": fb.like_count})


@app.route("/about")
def about():
    return render_template("about.html")


@app.route("/compiler")
@login_required
def compiler():
    return render_template("compiler.html")


@app.route("/aiassist")
@login_required
def aiassist():
    return render_template("aiassist.html")


@app.route("/register", methods=["GET", "POST"])
def register():
    if current_user.is_authenticated:
        return redirect(url_for("dashboard"))
    if request.method == "POST":
        username = request.form.get("username", "").strip()
        email = request.form.get("email", "").strip()
        password = request.form.get("password", "")
        confirm = request.form.get("confirm_password", "")

        if not username or not email or not password:
            flash("All fields are required.", "danger")
        elif password != confirm:
            flash("Passwords do not match.", "danger")
        elif password_errors(password):
            flash(
                "Password must contain: " + ", ".join(password_errors(password)) + ".",
                "danger",
            )
        elif User.query.filter_by(username=username).first():
            flash("That username is already taken.", "danger")
        elif User.query.filter_by(email=email).first():
            flash("An account with that email already exists.", "danger")
        else:
            user = User(
                username=username,
                email=email,
                password_hash=generate_password_hash(password),
            )
            db.session.add(user)
            db.session.commit()
            log_activity(user, "Registered an account", "New user signed up")
            flash("Account created successfully! You can now log in.", "success")
            return redirect(url_for("login", registered=1))
    return render_template("register.html")


@app.route("/login", methods=["GET", "POST"])
def login():
    if current_user.is_authenticated:
        return redirect(url_for("dashboard"))
    
    # Check if this might be the first login after database reset
    # Only show warning if no users exist in the database
    if User.query.count() == 0:
        show_database_reset_warning = True
    else:
        show_database_reset_warning = False
    
    if request.method == "POST":
        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")
        ip = _client_ip()
        user = User.query.filter_by(username=username).first()
        if is_ip_locked(ip) and not (user and user.is_admin):
            flash(
                "Too many failed login attempts from your IP address. "
                "You are blocked for about %d more minute%s." % (
                    ip_lock_remaining_minutes(ip),
                    "" if ip_lock_remaining_minutes(ip) == 1 else "s",
                ),
                "danger",
            )
            return render_template(
                "login.html", registered=request.args.get("registered") == "1", database_reset=show_database_reset_warning
            )
        if user and check_password_hash(user.password_hash, password):
            if user.is_banned:
                flash(
                    "Your account has been banned (%s). Contact an admin to get unbanned."
                    % (user.ban_reason or "violation of community rules"),
                    "danger",
                )
                return render_template(
                    "login.html", registered=request.args.get("registered") == "1", database_reset=show_database_reset_warning
                )
            login_user(user)
            log_activity(user, "Logged in", "Successful sign-in")
            flash(
                "Welcome back, " + user.username + "!",
                "success",
            )
            if user.is_admin:
                return redirect(url_for("admin_dashboard"))
            return redirect(url_for("dashboard"))
        record_failed_login(username)
        attempts_left = max(0, LOGIN_LOCK_ATTEMPTS - failed_login_count(ip))
        if attempts_left <= 0:
            flash(
                "Too many failed login attempts. Your IP address is now blocked "
                "from logging in for about %d more minute%s." % (
                    ip_lock_remaining_minutes(ip),
                    "" if ip_lock_remaining_minutes(ip) == 1 else "s",
                ),
                "danger",
            )
        else:
            flash(
                "Invalid username or password. "
                + ("1 attempt remaining" if attempts_left == 1 else "%d attempts remaining" % attempts_left)
                + " before your IP is blocked for %d minutes." % LOGIN_LOCK_MINUTES,
                "danger",
            )
    return render_template("login.html", registered=request.args.get("registered") == "1")


@app.route("/logout")
@login_required
def logout():
    log_activity(current_user, "Logged out", "User signed out")
    logout_user()
    flash("You have been logged out.", "info")
    return redirect(url_for("index"))


@app.route("/profile", methods=["POST"])
@login_required
def update_profile():
    username = request.form.get("username", "").strip()
    email = request.form.get("email", "").strip()
    new_password = request.form.get("new_password", "")
    confirm = request.form.get("confirm_password", "")
    avatar_file = request.files.get("avatar")
    remove_avatar = request.form.get("remove_avatar") == "yes"

    avatar_error = None
    avatar_data = None
    avatar_ext = ""
    if avatar_file and avatar_file.filename:
        avatar_ext = (
            avatar_file.filename.rsplit(".", 1)[-1].lower()
            if "." in avatar_file.filename
            else ""
        )
        if avatar_ext not in ALLOWED_AVATAR_EXT:
            avatar_error = "Only PNG, JPG, JPEG, GIF, or WebP images are allowed."
        else:
            avatar_data = avatar_file.read()
            if len(avatar_data) > 2 * 1024 * 1024:
                avatar_error = "Image must be 2 MB or smaller."

    if not username or not email:
        flash("Username and email are required.", "danger")
    elif User.query.filter(User.username == username, User.id != current_user.id).first():
        flash("That username is already taken.", "danger")
    elif User.query.filter(User.email == email, User.id != current_user.id).first():
        flash("An account with that email already exists.", "danger")
    elif new_password and new_password != confirm:
        flash("Passwords do not match.", "danger")
    elif new_password and password_errors(new_password):
        flash(
            "Password must contain: " + ", ".join(password_errors(new_password)) + ".",
            "danger",
        )
    elif avatar_error:
        flash(avatar_error, "danger")
    else:
        current_user.username = username
        current_user.email = email
        if new_password:
            current_user.password_hash = generate_password_hash(new_password)
        if avatar_data is not None:
            for old in glob.glob(os.path.join(AVATAR_DIR, f"avatar_{current_user.id}.*")):
                os.remove(old)
            mime = AVATAR_MIME.get(avatar_ext, "image/png")
            b64 = base64.b64encode(avatar_data).decode("ascii")
            current_user.avatar = f"data:{mime};base64,{b64}"
        elif remove_avatar:
            for old in glob.glob(os.path.join(AVATAR_DIR, f"avatar_{current_user.id}.*")):
                os.remove(old)
            current_user.avatar = None
        db.session.commit()
        log_activity(current_user, "Updated profile", "Changed username, email, avatar, or password")
        flash("Profile updated successfully.", "success")
    return redirect(url_for("dashboard"))


@app.route("/dashboard")
@login_required
def dashboard():
    completed = get_completed_ids(current_user)
    total = len(AVAILABLE_LESSONS)
    done = len(completed)
    percent = round((done / total) * 100) if total else 0
    scores = get_quiz_scores(current_user)
    quiz_points = sum(qs.best_score for qs in scores.values())
    quiz_count = len(scores)
    lab_ids = get_completed_lab_ids(current_user)
    lab_count = len(lab_ids)
    lab_points = lab_count * LAB_POINTS_PER_LAB
    total_points = quiz_points + lab_points
    sections = {}
    for lesson in LESSONS:
        sections.setdefault(lesson.get("section", "Fundamentals"), []).append(lesson)
    for name in SECTION_ORDER:
        sections.setdefault(name, [])
    locked_ids = set()
    unlocked_ids = get_unlocked_ids(current_user)
    for name in SECTION_ORDER:
        sec = sections.get(name, [])
        for i, l in enumerate(sec):
            if l.get("coming_soon"):
                continue
            if i > 0 and sec[i - 1]["id"] not in completed:
                if l["id"] not in unlocked_ids:
                    locked_ids.add(l["id"])
    section_quiz_map = {
        q.section: q
        for q in SectionQuizScore.query.filter_by(user_id=current_user.id).all()
    }
    all_quizzes_done = all(s in section_quiz_map for s in LONG_QUIZZES)
    section_stats = []
    for name in SECTION_ORDER:
        sec_lessons = [l for l in sections.get(name, []) if not l.get("coming_soon")]
        if not sec_lessons:
            continue
        sec_done = len([l for l in sec_lessons if l["id"] in completed])
        quiz_row = section_quiz_map.get(name)
        est_minutes = sum(lesson_estimate_minutes(l) for l in sec_lessons)
        section_stats.append(
            {
                "name": name,
                "total": len(sec_lessons),
                "done": sec_done,
                "percent": round((sec_done / len(sec_lessons)) * 100) if sec_lessons else 0,
                "quiz_best": quiz_row.best_score if quiz_row else None,
                "quiz_total": quiz_row.total_questions if quiz_row else None,
                "quiz_attempts": quiz_row.attempts if quiz_row else None,
                "estimate": est_minutes,
                "has_quiz": name in LONG_QUIZZES,
            }
        )
    lab_urls = {
        lid: (
            url_for(CUSTOM_LAB_ROUTES[lid])
            if lid in CUSTOM_LAB_ROUTES
            else url_for("lesson_lab", lesson_id=lid)
        )
        for lid in LAB_LESSON_IDS
    }
    return render_template(
        "dashboard.html",
        sections=sections,
        section_stats=section_stats,
        locked_ids=locked_ids,
        completed=completed,
        scores=scores,
        leaderboard=get_leaderboard(),
        total=total,
        done=done,
        percent=percent,
        quiz_points=quiz_points,
        quiz_count=quiz_count,
        total_points=total_points,
        lab_points=lab_points,
        lab_count=lab_count,
        all_lessons=AVAILABLE_LESSONS,
        section_minutes=SECTION_LESSON_MINUTES,
        lesson_estimate_minutes=lesson_estimate_minutes,
        lab_completed=bool(lab_count),
        completed_lab_ids=lab_ids,
        unlocked_ids=unlocked_ids,
        unlock_costs={l["id"]: lesson_unlock_cost(l["id"]) for l in LESSONS},
        lesson_difficulties={l["id"]: lesson_difficulty(l["id"]) for l in LESSONS},
        lab_lesson_ids=LAB_LESSON_IDS,
        lab_urls=lab_urls,
        all_quizzes_done=all_quizzes_done,
    )


@app.route("/admin")
@login_required
@admin_required
def admin_dashboard():
    now = datetime.utcnow()
    online_threshold = now - timedelta(minutes=5)
    total_users = User.query.count()
    online_users = User.query.filter(User.last_seen >= online_threshold).count()
    new_users_7d = User.query.filter(
        User.created_at >= now - timedelta(days=7)
    ).count()
    total_lessons_done = LessonProgress.query.count()
    total_quizzes = QuizScore.query.count()
    total_labs = LabProgress.query.count()
    total_logs = ActivityLog.query.count()
    failed_24h = FailedLogin.query.filter(
        FailedLogin.created_at >= now - timedelta(hours=24)
    ).count()
    recent_logs = (
        ActivityLog.query.order_by(ActivityLog.created_at.desc()).limit(15).all()
    )
    recent_users = User.query.order_by(User.created_at.desc()).limit(10).all()
    return render_template(
        "admin.html",
        total_users=total_users,
        online_users=online_users,
        new_users_7d=new_users_7d,
        total_lessons_done=total_lessons_done,
        total_quizzes=total_quizzes,
        total_labs=total_labs,
        total_logs=total_logs,
        failed_24h=failed_24h,
        recent_logs=recent_logs,
        recent_users=recent_users,
        now=now,
    )


@app.route("/admin/activity")
@login_required
@admin_required
def admin_activity():
    action = (request.args.get("action") or "").strip()
    try:
        page = max(1, int(request.args.get("page", 1)))
    except (TypeError, ValueError):
        page = 1
    per_page = 20
    q = ActivityLog.query
    if action:
        q = q.filter(ActivityLog.action == action)
    total = q.count()
    pages = max(1, (total + per_page - 1) // per_page)
    page = min(page, pages)
    logs = (
        q.order_by(ActivityLog.created_at.desc())
        .offset((page - 1) * per_page)
        .limit(per_page)
        .all()
    )
    actions = (
        db.session.query(ActivityLog.action)
        .distinct()
        .order_by(ActivityLog.action)
        .all()
    )
    return render_template(
        "admin_activity.html",
        logs=logs,
        actions=[a[0] for a in actions],
        current_action=action,
        log_count=total,
        page=page,
        pages=pages,
    )


@app.route("/admin/security")
@login_required
@admin_required
def admin_security():
    now = datetime.utcnow()
    since_24h = now - timedelta(hours=24)
    total_24h = FailedLogin.query.filter(
        FailedLogin.created_at >= since_24h
    ).count()
    total_all = FailedLogin.query.count()
    ip_rows = {}
    for fl in FailedLogin.query.order_by(FailedLogin.created_at.desc()).all():
        if not fl.ip_address:
            continue
        entry = ip_rows.setdefault(
            fl.ip_address,
            {
                "ip": fl.ip_address,
                "count": 0,
                "last_username": None,
                "last_seen": None,
                "locked": False,
            },
        )
        entry["count"] += 1
        if entry["last_username"] is None and fl.username:
            entry["last_username"] = fl.username
        if entry["last_seen"] is None or fl.created_at > entry["last_seen"]:
            entry["last_seen"] = fl.created_at
    for entry in ip_rows.values():
        entry["locked"] = failed_login_count(entry["ip"]) >= LOGIN_LOCK_ATTEMPTS
    locked_count = sum(1 for e in ip_rows.values() if e["locked"])
    recent = (
        FailedLogin.query.order_by(FailedLogin.created_at.desc()).limit(100).all()
    )
    return render_template(
        "admin_security.html",
        total_24h=total_24h,
        total_all=total_all,
        locked_count=locked_count,
        ip_rows=sorted(
            ip_rows.values(), key=lambda e: e["count"], reverse=True
        ),
        recent=recent,
        lock_attempts=LOGIN_LOCK_ATTEMPTS,
        lock_minutes=LOGIN_LOCK_MINUTES,
        now=now,
    )


@app.route("/admin/security/clear", methods=["POST"])
@login_required
@admin_required
def admin_security_clear():
    ip = (request.form.get("ip") or "").strip()
    if ip:
        clear_failed_logins(ip)
        flash("Cleared failed login attempts for %s." % ip, "success")
    else:
        clear_failed_logins()
        flash("Cleared all failed login attempts.", "success")
    return redirect(url_for("admin_security"))


@app.route("/admin/feedback")
@login_required
@admin_required
def admin_feedback():
    feedbacks = Feedback.query.order_by(Feedback.created_at.desc()).all()
    unread_count = Feedback.query.filter_by(is_read=False).count()
    return render_template("admin_feedback.html", feedbacks=feedbacks, unread_count=unread_count)


@app.route("/admin/feedback/<int:fid>/read", methods=["POST"])
@login_required
@admin_required
def admin_feedback_read(fid):
    fb = db.session.get(Feedback, fid)
    if fb:
        fb.is_read = True
        db.session.commit()
    flash("Feedback marked as read.", "success")
    return redirect(url_for("admin_feedback"))


@app.route("/admin/feedback/<int:fid>/reply", methods=["POST"])
@login_required
@admin_required
def admin_feedback_reply(fid):
    fb = db.session.get(Feedback, fid)
    if not fb:
        flash("Feedback not found.", "danger")
        return redirect(url_for("admin_feedback"))
    reply = request.form.get("reply", "").strip()
    if reply:
        fb.admin_reply = reply
        fb.replied_at = datetime.utcnow()
        fb.is_read = True
        db.session.commit()
        log_activity(current_user, "Replied to feedback", f"Feedback #{fid} from {fb.name}")
        flash("Reply sent successfully.", "success")
    return redirect(url_for("admin_feedback"))


@app.route("/admin/feedback/<int:fid>/delete", methods=["POST"])
@login_required
@admin_required
def admin_feedback_delete(fid):
    fb = db.session.get(Feedback, fid)
    if fb:
        FeedbackLike.query.filter_by(feedback_id=fid).delete()
        db.session.delete(fb)
        db.session.commit()
    flash("Feedback deleted.", "success")
    return redirect(url_for("admin_feedback"))


@app.route("/admin/users")
@login_required
@admin_required
def admin_users():
    now = datetime.utcnow()
    online_threshold = now - timedelta(minutes=5)
    users = User.query.order_by(User.created_at.desc()).all()
    rows = []
    online_count = 0
    for u in users:
        is_online = u.last_seen and u.last_seen >= online_threshold
        if is_online:
            online_count += 1
        rows.append(
            {
                "user": u,
                "lessons_done": len(u.progress),
                "quizzes": len(u.quiz_scores),
                "labs": len(db.session.query(LabProgress).filter_by(user_id=u.id).all()),
                "is_online": is_online,
                "last_seen": u.last_seen,
            }
        )
    return render_template("admin_users.html", rows=rows, total=len(rows), total_lessons=len(AVAILABLE_LESSONS), online_count=online_count)


@app.route("/admin/user/<int:uid>/ban", methods=["POST"])
@login_required
@admin_required
def admin_ban_user(uid):
    user = db.session.get(User, uid)
    if user is None:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))
    if user.is_admin:
        flash("You cannot ban an admin account.", "danger")
        return redirect(url_for("admin_users"))
    reason = (request.form.get("reason") or "").strip() or "Violation of community rules"
    if ban_user(user, reason):
        flash("Banned %s." % user.username, "success")
    else:
        flash("%s is already banned." % user.username, "warning")
    return redirect(url_for("admin_users"))


@app.route("/admin/user/<int:uid>/unban", methods=["POST"])
@login_required
@admin_required
def admin_unban_user(uid):
    user = db.session.get(User, uid)
    if user is None:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))
    if unban_user(user):
        flash("Unbanned %s." % user.username, "success")
    else:
        flash("%s is not banned." % user.username, "warning")
    return redirect(url_for("admin_users"))


@app.route("/admin/user/<int:uid>/coins", methods=["POST"])
@login_required
@admin_required
def admin_set_coins(uid):
    user = db.session.get(User, uid)
    if user is None:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))
    raw = (request.form.get("coins") or "").strip()
    try:
        amount = int(raw)
    except ValueError:
        flash("Please enter a valid whole number of coins.", "danger")
        return redirect(url_for("admin_users"))
    if amount < 0:
        flash("Coins cannot be negative.", "danger")
        return redirect(url_for("admin_users"))
    old = user.coins or 0
    user.coins = amount
    db.session.commit()
    log_activity(
        current_user,
        "Admin set coins",
        "Set %s's coins from %d to %d" % (user.username, old, amount),
    )
    flash("Set %s's coins to %d." % (user.username, amount), "success")
    return redirect(url_for("admin_users"))


@app.route("/admin/api/online-users")
@login_required
@admin_required
def api_online_users():
    now = datetime.utcnow()
    online_threshold = now - timedelta(minutes=5)
    users = User.query.order_by(User.created_at.desc()).all()
    rows = []
    online_count = 0
    for u in users:
        is_online = u.last_seen and u.last_seen >= online_threshold
        if is_online:
            online_count += 1
        rows.append({
            "id": u.id,
            "username": u.username,
            "email": u.email,
            "avatar": u.avatar,
            "is_admin": u.is_admin,
            "is_online": is_online,
            "last_seen": u.last_seen.strftime("%b %d, %I:%M %p") if u.last_seen else "Never",
            "lessons_done": len(u.progress),
            "quizzes": len(u.quiz_scores),
            "labs": len(db.session.query(LabProgress).filter_by(user_id=u.id).all()),
        })
    return jsonify({"rows": rows, "total": len(rows), "online_count": online_count})


@app.route("/admin/api/stats")
@login_required
@admin_required
def api_admin_stats():
    now = datetime.utcnow()
    online_threshold = now - timedelta(minutes=5)
    return jsonify({
        "total_users": User.query.count(),
        "online_users": User.query.filter(User.last_seen >= online_threshold).count(),
        "new_users_7d": User.query.filter(User.created_at >= now - timedelta(days=7)).count(),
        "total_lessons_done": LessonProgress.query.count(),
        "total_quizzes": QuizScore.query.count(),
        "total_labs": LabProgress.query.count(),
    })


@app.route("/admin/user/<int:uid>/progress")
@login_required
@admin_required
def admin_user_progress(uid):
    user = db.session.get(User, uid)
    if not user:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))

    completed = get_completed_ids(user)
    scores = get_quiz_scores(user)
    lab_ids = get_completed_lab_ids(user)

    # Build lesson progress list
    lessons_progress = []
    for l in LESSONS:
        if l.get("coming_soon"):
            continue
        lid = l["id"]
        lessons_progress.append({
            "lesson": l,
            "completed": lid in completed,
            "quiz_score": scores.get(lid),
            "lab_completed": lid in lab_ids,
        })

    # Section quiz scores
    section_scores = {}
    for name in SECTION_ORDER:
        sqs = SectionQuizScore.query.filter_by(user_id=user.id, section=name).first()
        section_scores[name] = sqs

    return render_template(
        "admin_user_progress.html",
        target_user=user,
        lessons_progress=lessons_progress,
        section_scores=section_scores,
        total_lessons=len([l for l in LESSONS if not l.get("coming_soon")]),
    )


@app.route("/admin/user/<int:uid>/reset-lesson/<int:lid>", methods=["POST"])
@login_required
@admin_required
def admin_reset_lesson(uid, lid):
    user = db.session.get(User, uid)
    if not user:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))

    lesson = next((l for l in LESSONS if l["id"] == lid), None)
    if not lesson:
        flash("Lesson not found.", "danger")
        return redirect(url_for("admin_user_progress", uid=uid))

    # Remove lesson completion
    deleted = LessonProgress.query.filter_by(user_id=uid, lesson_id=lid).delete()
    # Remove quiz score
    deleted_quiz = QuizScore.query.filter_by(user_id=uid, lesson_id=lid).delete()
    # Remove lab completion if applicable
    deleted_lab = LabProgress.query.filter_by(user_id=uid, lesson_id=lid).delete()
    user.coins = len(LessonProgress.query.filter_by(user_id=uid).all()) * COINS_PER_LESSON
    db.session.commit()

    log_activity(user, "Admin reset lesson", f"Admin reset lesson '{lesson['title']}' (completed={deleted}, quiz={deleted_quiz}, lab={deleted_lab})")
    flash(f"Reset lesson '{lesson['title']}' for {user.username}.", "success")
    return redirect(url_for("admin_user_progress", uid=uid))


@app.route("/admin/user/<int:uid>/reset-section-quiz/<section>", methods=["POST"])
@login_required
@admin_required
def admin_reset_section_quiz(uid, section):
    user = db.session.get(User, uid)
    if not user:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))

    deleted = SectionQuizScore.query.filter_by(user_id=uid, section=section).delete()
    db.session.commit()

    log_activity(user, "Admin reset section quiz", f"Admin reset section quiz for '{section}' (deleted={deleted})")
    flash(f"Reset section quiz '{section}' for {user.username}.", "success")
    return redirect(url_for("admin_user_progress", uid=uid))


@app.route("/admin/user/<int:uid>/reset-all", methods=["POST"])
@login_required
@admin_required
def admin_reset_all(uid):
    user = db.session.get(User, uid)
    if not user:
        flash("User not found.", "danger")
        return redirect(url_for("admin_users"))

    deleted_progress = LessonProgress.query.filter_by(user_id=uid).delete()
    deleted_quizzes = QuizScore.query.filter_by(user_id=uid).delete()
    deleted_section_quizzes = SectionQuizScore.query.filter_by(user_id=uid).delete()
    deleted_labs = LabProgress.query.filter_by(user_id=uid).delete()
    user.coins = 0
    db.session.commit()

    log_activity(user, "Admin reset all progress", f"Admin reset ALL progress (lessons={deleted_progress}, quizzes={deleted_quizzes}, section_quizzes={deleted_section_quizzes}, labs={deleted_labs})")
    flash(f"Reset ALL progress for {user.username}.", "success")
    return redirect(url_for("admin_user_progress", uid=uid))


@app.route("/certificate/<path:name>")
@login_required
def certificate(name):
    completed = set(get_completed_ids(current_user))
    quiz_map = {
        q.section: q
        for q in SectionQuizScore.query.filter_by(user_id=current_user.id).all()
    }
    if name == "Full Course":
        if len(completed) < len(AVAILABLE_LESSONS):
            flash(
                f"Complete all {len(AVAILABLE_LESSONS)} lessons to unlock your full course certificate.",
                "warning",
            )
            return redirect(url_for("dashboard"))
        missing_quizzes = [s for s in LONG_QUIZZES if s not in quiz_map]
        if missing_quizzes:
            flash(
                "Finish every long quiz to unlock your full course certificate: "
                + ", ".join(missing_quizzes)
                + ".",
                "warning",
            )
            return redirect(url_for("dashboard"))
        cert_title = "Course Completion"
        cert_subtitle = "Complete Curriculum - Programming Fundamentals and Advanced Database Systems"
        sec_num = 0
    else:
        sec_lessons = [
            l
            for l in LESSONS
            if l.get("section", "Fundamentals") == name and not l.get("coming_soon")
        ]
        if not sec_lessons:
            flash("Certificate not found.", "danger")
            return redirect(url_for("dashboard"))
        done_count = len([l for l in sec_lessons if l["id"] in completed])
        if done_count < len(sec_lessons):
            flash(f"Complete all {len(sec_lessons)} lessons in the {name} section to unlock this certificate.", "warning")
            return redirect(url_for("dashboard"))
        if name in LONG_QUIZZES and name not in quiz_map:
            flash(
                f"Finish the {name} long quiz to unlock this certificate.",
                "warning",
            )
            return redirect(url_for("dashboard"))
        cert_title = f"{name} Completion"
        cert_subtitle = f"All {len(sec_lessons)} lessons in the {name} section"
        sec_num = (SECTION_ORDER.index(name) + 1) if name in SECTION_ORDER else 0
    issued = datetime.utcnow()
    cert_no = f"CF-{issued.year}-{current_user.id:03d}{sec_num:02d}"
    return render_template(
        "certificate.html",
        cert_no=cert_no,
        cert_title=cert_title,
        cert_subtitle=cert_subtitle,
        section=name,
        issued=issued,
    )


def _gemini_api_key():
    return os.environ.get("GEMINI_API_KEY", "").strip()


def _gemini_model():
    return os.environ.get("GEMINI_MODEL", "gemini-3.5-flash-lite").strip()

EREN_SYSTEM_PROMPT = (
    "You are E.R.E.N (Educational Response Engine for Novices), the friendly AI assistant of "
    "CodeFundamentals, a free site where students learn Python, C++, Java, and C side by side, plus "
    "Advanced Database System (SQL). ALWAYS respond only in English or Filipino (Tagalog) — "
    "never in any other language. If the student writes in English or Filipino, answer in that "
    "same language. If the student writes in any other language (for example Spanish, French, "
    "German, or Cebuano), still answer in English or Filipino — default to English. Be concise, "
    "patient, and beginner-friendly. Answer questions about the course, its lessons, quizzes, "
    "leaderboard, and coding concepts. Never refuse to answer — you can answer questions on any "
    "topic, even off-topic ones like general knowledge or personal help, and if a question is "
    "about the course, use it as a chance to teach.\n\n"
    "Course sections and lessons:\n"
    "- Fundamentals: Introduction to Programming, Programming History, Variables and Data Types, "
    "Input and Output, Operators, Logical Operators, Conditionals (if/else), Loops, Functions, "
    "Arrays and Lists, Basic Problem Solving. C is taught side by side with Python, C++, and Java "
    "inside every Fundamentals lesson.\n"
    "- Advanced Fundamentals: Introduction to GitHub, Web Deployment Fundamentals, Deploying with Render, Deploying with Netlify, Building Websites with WordPress.\n"
    "- Advanced Database System: DBMS overview, ER model and keys, Normalization (1NF-3NF), "
    "SQL DDL, DML, SELECT and filtering, JOINs, Subqueries and Aggregates.\n"
    "- SQL Injection: Introduction to SQL Injection, How SQL Injection Works, Types of SQL Injection, "
    "UNION-Based SQL Injection, Blind SQL Injection, SQL Injection in Real-World Queries, "
    "Preventing SQL Injection: Parameterized Queries, Defense in Depth, SQL Injection Testing and Detection, "
    "and Real-World Attack Scenarios. Every SQL Injection lesson has its own interactive Practice Lab "
    "(2 points each, marked green on the dashboard).\n\n"
    "Rules: each lesson ends with a quiz (one attempt). Each section ends with a Long Quiz "
    "(15 questions, 20 seconds each, 2 attempts). The leaderboard ranks users by total points: "
    "quiz scores plus 2 points per SQL Injection lab.\n\n"
    "Video Lessons & Embedding:\n"
    "When a student asks to watch, see, or receive a video lesson for a topic (e.g., 'papanood ng video', 'send video lesson', 'video on loops', etc.), you CAN and SHOULD embed the video using YouTube iframe format: `<iframe src=\"https://www.youtube.com/embed/VIDEO_ID\" title=\"TITLE\" allowfullscreen></iframe>` or `https://www.youtube.com/embed/VIDEO_ID`.\n"
    "Available lesson video IDs:\n"
    "- Introduction to Programming: Python (9QKe2RvjG-A), C++ (LApkQYRUru8), Java (Pyx9oLYpbi4), C (aMpsKnf6DrQ)\n"
    "- Programming History: General (RU1u-js7db8), Python (GfH4QL4VqJ0), C++ (lI7tMxzSJ7w), Java (ZqGSg4b_cZA), C (6VT8hDr2GhU)\n"
    "- Variables & Data Types: Python (2SG133xkpMQ), C++ (E4yj6QX-TW8), Java (cD-6C39RPKs), C (OSyjOvFbAGI)\n"
    "- Input & Output: Python (7iZOEoGWnuU), C++ (vkALQRM5NNU), Java (sk_TzAjZ3zs), C (xOIVXR35aI4)\n"
    "- Operators: Python (RYEDDygFNx8), C++ (ib-Qfx7-iIQ), Java (M0lnSma7qMw), C (vvpDbhqPrww)\n"
    "- Logical Operators: Python (srYNLwnEfeo), C++ (JrcSEPKn6uI), Java (h7pGky1QaJs), C (U19kiynYopE)\n"
    "- Conditionals: Python (aXWb4rDK8NE), C++ (ISLTa95qJJA), Java (GaAKcy9cEcU), C (q2LCT6gRZVY)\n"
    "- Loops: Python (EX47v75YM1Y), C++ (-Qev7T2a_Lc), Java (Mt-gZKX3dq0), C (b4DPj0XAfSg)\n"
    "- Functions: Python (89cGQjB5R4M), C++ (67I3ZEmyVKQ), Java (v5p_SUfi710), C (NGQoKF2Ggt8)\n"
    "- Arrays & Lists: Python (dpwp1NADaFU), C++ (dRyf8cLyKgA), Java (9dr2mHYYoug), C (MOeGnamlUP4)\n"
    "- Basic Problem Solving: (6XJ8294lC0c)\n"
    "- Database Systems Overview: (wR0jg0eQsZA)\n"
    "- ER Model & Keys: (xsg9BDiwiJE)\n"
    "- Normalization: (GFQaEYEc8_8)\n"
    "- SQL Basics & DDL: (3Qbq61A_sNs)\n"
    "- SQL DML: (ku3vMAP0h0s)\n"
    "- SQL SELECT & Filtering: (4Uv0o8IBqw0)\n"
    "- SQL JOINs: (Yh4CrPHVBdE)\n"
    "- Subqueries & Aggregates: (GpC0XyiJPEo)\n"
    "- Web Deployment Fundamentals: Laravel (m-QO9Qp_wRQ), React (SuZBpX7Y7EA)\n"
    "- Introduction to GitHub: (tRZGeaHPoaw)\n"
    "- Deploying with Render: (srudWEWKPv0)\n"
    "- Deploying with Netlify: (IH95RG9Y6L4, YUtNzwxOVYY)\n"
    "- SQL Injection Intro: (wcaiKgQU6VE)\n"
    "- How SQL Injection Works: (FwIUkAwKzG8)\n"
    "- Angular (Advanced Fundamentals): (cRC9DlH45lA)\n"
    "\n\nYouTube Summaries:\n"
    "When a student pastes a YouTube link, the system automatically fetches the video transcript "
    "and asks you to summarize it. Give a beginner-friendly summary (English or Filipino) with "
    "short ## headers and - bullets, highlight key points, and answer any question they asked.\n"
    "\n\nPresentations / PPT requests:\n"
    "When a student asks for a presentation, PowerPoint, slide outline, or to 'make this into slides' for any "
    "topic, produce a clean slide deck using ONLY this format — it renders nicely in chat and can be exported "
    "to a real .pptx file:\n"
    "- Start with a title slide heading: `## Slide 1: <Title>`\n"
    "- Each following slide heading starts with `##` and the slide number, e.g. `## Slide 2: <Topic>`.\n"
    "- Under each heading add short `-` bullet points (one level, no nested bullets).\n"
    "- Do NOT use `**`, `*`, `#` (beyond the `##` heading), tables, or other markdown in the outline.\n"
    "- You may add one short code example per slide inside a triple-backtick ``` fence; it is shown in chat "
    "and automatically included in the exported PowerPoint.\n"
    "- After every code example, always include its sample output on its own line as \"Sample Output:\" "
    "followed by a ```output fence containing the exact printed result, e.g.:\n"
    "```c\nprintf(\"Hello\");\n```\nSample Output:\n```output\nHello\n```\n"
    "- Keep it to 6-10 slides with concise bullets, and no closing message after the outline."
)


def gemini_reply(message, history=None):
    api_key = _gemini_api_key()
    if not api_key:
        print("[E.R.E.N] No GEMINI_API_KEY set, using offline rules.", file=sys.stderr)
        return None
    model = _gemini_model()
    contents = []
    for item in (history or [])[-14:]:
        role = item.get("role")
        content = (item.get("content") or "").strip()
        if role not in ("user", "assistant") or not content:
            continue
        contents.append(
            {
                "role": "user" if role == "user" else "model",
                "parts": [{"text": content}],
            }
        )
    contents.append({"role": "user", "parts": [{"text": message}]})
    cleaned = []
    for entry in contents:
        if cleaned and cleaned[-1]["role"] == entry["role"]:
            cleaned[-1]["parts"][0]["text"] += "\n\n" + entry["parts"][0]["text"]
        else:
            cleaned.append(entry)
    while cleaned and cleaned[0]["role"] != "user":
        cleaned.pop(0)
    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        f"{model}:generateContent?key={api_key}"
    )
    payload = {
        "system_instruction": {"parts": [{"text": EREN_SYSTEM_PROMPT}]},
        "contents": cleaned,
        "generationConfig": {"temperature": 0.7, "maxOutputTokens": 4096},
    }
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=25) as resp:
            body = json.loads(resp.read().decode("utf-8"))
        parts = body.get("candidates", [{}])[0].get("content", {}).get("parts") or []
        reply = "".join(p.get("text", "") for p in parts).strip()
        return reply or None
    except Exception as exc:
        print(f"[E.R.E.N] Gemini request failed ({model}), using offline rules: {exc}", file=sys.stderr)
        return None


GENERAL_KNOWLEDGE = [
    (
        ["research", "pananaliksik"],
        "Research is the systematic process of gathering, analyzing, and interpreting "
        "information to answer a question or solve a problem. The usual steps are: identify "
        "the problem, review what is already known, collect data, analyze it, and draw "
        "conclusions.",
    ),
    (
        ["algorithm", "algorithms", "algoritmo"],
        "An algorithm is a step-by-step set of instructions for solving a problem, like a "
        "recipe. Planning with pseudocode is designing an algorithm before you translate it "
        "into Python, C++, or Java.",
    ),
    (
        ["data structure", "data structures"],
        "A data structure is a way of organizing and storing data so your program can use it "
        "efficiently. Arrays, lists, stacks, and queues are examples — choosing the right one "
        "makes your code faster and cleaner.",
    ),
    (
        ["artificial intelligence", "ai", "a.i."],
        "Artificial Intelligence (AI) is the field that builds programs that can learn, "
        "reason, and make decisions. Chatbots, face unlock, and recommendation engines all "
        "use AI — even E.R.E.N has a little AI brain!",
    ),
    (
        ["machine learning"],
        "Machine learning is a branch of AI where a program learns patterns from data instead "
        "of following fixed rules. Spam filters, voice assistants, and Netflix suggestions "
        "are everyday examples.",
    ),
    (
        ["computer science"],
        "Computer science is the study of computers and computation: how data is represented, "
        "how algorithms work, and how software is built. Writing code is one important part, "
        "but not the only one.",
    ),
    (
        ["software"],
        "Software is the set of instructions that tells a computer what to do — the programs "
        "themselves. Your operating system, your web browser, and CodeFundamentals are all "
        "software.",
    ),
    (
        ["hardware"],
        "Hardware is the physical parts of a computer you can touch: the CPU, RAM, hard "
        "drive, keyboard, and monitor. Software runs on top of hardware.",
    ),
    (
        ["operating system", "os"],
        "An operating system (OS) is the core software that manages your computer's hardware "
        "and runs other programs. Windows, macOS, and Linux are operating systems.",
    ),
    (
        ["browser", "browsers"],
        "A web browser is the program you use to view websites. Chrome, Firefox, Edge, and "
        "Safari are browsers — you are using one right now to reach CodeFundamentals.",
    ),
    (
        ["internet"],
        "The internet is a global network of connected computers that share information. "
        "Websites, email, online games, and video calls all travel over the internet.",
    ),
    (
        ["network", "networks", "networking"],
        "A network is a group of connected computers that can share data and resources. The "
        "internet is the largest network in the world.",
    ),
    (
        ["server", "servers"],
        "A server is a computer that stores websites and data and sends them to other "
        "computers (clients) when requested. When you open a webpage, your browser asks a "
        "server for it.",
    ),
    (
        ["website", "websites", "web page", "web pages", "webpage", "webpages", "web"],
        "A website is a collection of connected pages available on the internet. It is "
        "usually built with HTML, CSS, and JavaScript on the front end, and something like "
        "Python or SQL on the back end — just like this site.",
    ),
    (
        ["application", "applications", "app", "apps"],
        "An application (app) is a program built for a specific task, like a calculator, a "
        "messaging app, or a browser. Mobile apps and desktop apps are both applications.",
    ),
    (
        ["git", "github", "version control"],
        "Git is a tool that tracks changes to your code over time — that is version control. "
        "GitHub is a website that stores Git repositories online so you can share and "
        "collaborate on code with others.",
    ),
    (
        ["ide", "editor"],
        "An IDE (Integrated Development Environment) is a program that helps you write code "
        "with features like syntax highlighting, auto-complete, and a built-in run button — "
        "for example VS Code, PyCharm, and IntelliJ.",
    ),
    (
        ["career", "job", "trabaho"],
        "A programming career can mean many things: web developer, game developer, data "
        "scientist, or systems engineer. They all start with the same basics you are learning "
        "here — variables, loops, functions, and arrays.",
    ),
    (
        ["study", "studying", "review", "revising", "mag-aral", "aral", "tips"],
        "Study tips: code a little every day, explain concepts out loud (or to E.R.E.N!), "
        "practice with the code editor, and retake quizzes until the ideas stick. Ten focused "
        "minutes beats an hour of distraction.",
    ),
    (
        ["motivation", "motivated", "motivate", "give up", "too hard", "struggling"],
        "Everyone struggles with programming sometimes — even professionals! Break the problem "
        "into tiny pieces, fix one bug at a time, and remember why you started. Slow progress "
        "is still progress.",
    ),
    (
        ["time management", "procrastinate", "procrastination", "distracted", "focus"],
        "To beat procrastination: pick one small task, set a 25-minute timer, and silence "
        "your phone. Studying in short focused bursts with short breaks works better than "
        "long, distracted sessions.",
    ),
    (
        ["mathematics", "math", "matematika"],
        "Math is the study of numbers, patterns, and relationships. Programming leans on "
        "logic and arithmetic — operators like +, -, *, /, and % are math you already use in "
        "the lessons.",
    ),
    (
        ["science", "agham"],
        "Science is the study of the natural world through observation, testing, and evidence "
        "— the scientific method. Computer science is the branch that studies computation and "
        "software.",
    ),
    (
        ["technology", "tech"],
        "Technology is the use of scientific knowledge to build tools that make life easier — "
        "from the smartphone in your pocket to the servers that run this website.",
    ),
    (
        ["history", "kasaysayan"],
        "The history of programming is a whole lesson in the Fundamentals section! It traces "
        "how computers and languages like Python, C++, and Java evolved, and even includes "
        "documentaries you can watch right in the lesson.",
    ),
    (
        ["deploy", "deployment", "deploying", "hosting", "production"],
        "Deployment is the process of moving your app from your local machine to a server "
        "so anyone on the internet can use it. The Fundamentals section has a full lesson on "
        "Web Deployment Fundamentals covering Laravel, React, hosting platforms, and CI/CD. "
        "Ask me for a video lesson on deployment!",
    ),
    (
        ["domain", "domain name", "dns"],
        "A domain name (like example.com) is a human-readable address that points to your "
        "server's IP address. DNS (Domain Name System) translates the domain into the actual "
        "IP. You can buy domains from registrars like Namecheap or Cloudflare and point them "
        "to your hosting provider.",
    ),
    (
        ["ci/cd", "ci cd", "continuous integration", "continuous deployment"],
        "CI/CD (Continuous Integration / Continuous Deployment) is a practice where code "
        "changes are automatically tested and deployed. GitHub Actions is a popular tool "
        "that runs your tests and deploys your app every time you push to GitHub.",
    ),
    (
        ["vercel", "netlify", "github pages", "cloudflare pages"],
        "Vercel, Netlify, GitHub Pages, and Cloudflare Pages are free hosting platforms "
        "perfect for deploying React apps and static sites. They connect to your GitHub "
        "repo and auto-deploy on every push. The Web Deployment Fundamentals lesson covers this!",
    ),
    (
        ["github", "git", "repository", "repo", "commit", "branch", "merge", "pull request"],
        "GitHub is the world's largest code hosting platform. Git is the version control "
        "tool that tracks changes. The Fundamentals section has a full lesson on "
        "Introduction to GitHub covering repos, commits, branches, and pull requests. "
        "Ask me for a video lesson on GitHub!",
    ),
    (
        ["render", "render.com", "render deploy", "deploy on render"],
        "Render is a cloud platform (PaaS) that makes deploying web apps effortless. "
        "Connect your GitHub repo, set your build/start commands, and Render handles SSL, "
        "auto-deploys, and scaling. The Deploying with Render lesson covers Laravel, React, "
        "and Node.js deployments with free tier support. Ask me for a video lesson!",
    ),
    (
        ["netlify", "netlify.com", "netlify deploy", "deploy on netlify"],
        "Netlify is a web platform for building, deploying, and managing modern web projects. "
        "Deploy static sites, SPAs, and JAMstack apps instantly from Git. Netlify offers free "
        "SSL, drag-and-drop deployment, serverless functions, and form handling. The Deploying "
        "with Netlify lesson covers React, Vue, and HTML/CSS/JS deployments. Ask me for a video lesson!",
    ),
    (
        ["version control", "version control system"],
        "Version control tracks every change to your code over time. Git is the most "
        "popular version control system, and GitHub is where you store Git repositories "
        "online. This lets you undo mistakes, collaborate with teammates, and deploy with "
        "CI/CD. Check out the Introduction to GitHub lesson!",
    ),
]


def general_knowledge_reply(text):
    for keywords, answer in GENERAL_KNOWLEDGE:
        if any(re.search(r"(?<!\w)" + re.escape(k) + r"(?!\w)", text) for k in keywords):
            return answer
    return None


def is_definition_question(text):
    return bool(
        re.search(
            r"^(?:what\s+is\s+(?:a|an|the)\s+|what\s+is\s+|what\s+are\s+|"
            r"what's\s+|whats\s+|define\s+|meaning\s+of\s+|explain\s+|"
            r"ano\s+ba\s+ang\s+|ano\s+ang\s+|ano\s+yung\s+|ano\s+yang\s+|ano\s+)",
            text,
        )
    )


def get_video_reply(text):
    msg_lower = text.lower()
    video_triggers = [
        "video", "papanood", "panuorin", "panoorin", "papanuod", "panood",
        "watch", "embed", "iframe", "papanoodin", "pasend ng video", "pa send ng video",
        "patingin ng video", "tutorial video", "video lesson", "send video"
    ]
    is_asking_video = any(vt in msg_lower for vt in video_triggers)
    if not is_asking_video:
        return None

    lang_filter = None
    if "python" in msg_lower:
        lang_filter = "python"
    elif "c++" in msg_lower or "cpp" in msg_lower or "c plus plus" in msg_lower:
        lang_filter = "cpp"
    elif "java" in msg_lower:
        lang_filter = "java"

    topic_mappings = [
        (["variable", "variables", "data type", "datatype", "integer", "boolean", "float", "string", "char"], 2),
        (["input", "output", "io", "print", "scanner", "cin", "cout"], 3),
        (["operator", "operators", "arithmetic", "modulo", "remainder"], 4),
        (["logical", "logical operator", "and", "or", "not"], 5),
        (["conditional", "conditionals", "if else", "if-else", "elif", "else if", "decision"], 6),
        (["loop", "loops", "for loop", "while", "repeat", "iteration"], 7),
        (["function", "functions", "method", "methods", "def", "return"], 8),
        (["array", "arrays", "list", "lists", "vector", "collection"], 9),
        (["problem solving", "pseudocode", "flowchart", "algorithm"], 10),
        (["history", "kasaysayan", "creator", "story"], 29),
        (["dbms", "database overview", "database system", "database"], 11),
        (["er model", "erd", "entity", "key", "primary key", "foreign key"], 12),
        (["normalization", "1nf", "2nf", "3nf"], 13),
        (["ddl", "create table", "alter", "drop"], 14),
        (["dml", "insert", "update", "delete"], 15),
        (["select", "where", "filter", "dql"], 16),
        (["join", "joins", "inner join", "left join"], 17),
        (["subquery", "subqueries", "aggregate", "group by", "having"], 18),
        (["deploy", "deployment", "deploying", "hosting", "host", "production", "ci/cd", "ci cd", "vercel", "github pages", "laravel deploy", "react deploy"], 30),
        (["github", "git", "repository", "repo", "commit", "branch", "merge", "pull request", "clone"], 31),
        (["render", "deploy render", "render deploy", "render.com", "render yaml"], 32),
        (["netlify", "deploy netlify", "netlify deploy", "netlify.com", "netlify.toml", "netlify drop"], 33),
        (["sql injection intro", "sqli intro", "what is sql injection"], 19),
        (["sql injection", "sqli", "injection"], 20),
        (["intro", "introduction", "basics", "programming"], 1),
    ]

    matched_lesson = None
    for keywords, lid in topic_mappings:
        if any(re.search(r"(?<!\w)" + re.escape(k) + r"(?!\w)", msg_lower) for k in keywords):
            matched_lesson = next((l for l in LESSONS if l["id"] == lid), None)
            if matched_lesson:
                break

    if not matched_lesson:
        matched_lesson = next((l for l in LESSONS if l["id"] == 1), None)

    if not matched_lesson:
        return "I couldn't find a video lesson for that topic yet. Try asking for a video lesson on **Variables**, **Operators**, **Conditionals**, **Loops**, **Functions**, **Arrays**, or **SQL**!"

    videos_data = matched_lesson.get("videos")
    if not videos_data:
        return f"No video lesson is currently available for **{matched_lesson['title']}**."

    video_items = []
    if isinstance(videos_data, dict):
        if lang_filter and lang_filter in videos_data and videos_data[lang_filter]:
            video_items = videos_data[lang_filter]
        elif "general" in videos_data and videos_data["general"]:
            video_items = videos_data["general"]
        else:
            for lk in ["python", "cpp", "java"]:
                if lk in videos_data and videos_data[lk]:
                    video_items.extend(videos_data[lk])
                    break
    elif isinstance(videos_data, list):
        video_items = videos_data

    if not video_items:
        return f"No video lesson found for **{matched_lesson['title']}**."

    reply_lines = [f"Here is the video lesson for **{matched_lesson['title']}**:\n"]
    for v in video_items[:2]:
        v_title = v.get("title", matched_lesson["title"])
        v_id = v.get("id")
        v_channel = v.get("channel", "")
        channel_str = f" ({v_channel})" if v_channel else ""
        reply_lines.append(f"**{v_title}**{channel_str}")
        reply_lines.append(f'<iframe src="https://www.youtube.com/embed/{v_id}" title="{v_title}" allowfullscreen></iframe>\n')

    reply_lines.append("You can play and watch the video lesson directly right here! Let me know if you need another video lesson or topic.")
    return "\n".join(reply_lines)


def assistant_reply(message):
    text = message.lower().strip()
    v_reply = get_video_reply(text)
    if v_reply:
        return v_reply

    def brief(lesson_id, lead):
        lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
        if not lesson:
            return lead
        section_lessons = [
            l for l in LESSONS
            if l.get("section", "Fundamentals") == lesson.get("section", "Fundamentals")
        ]
        num = next(i + 1 for i, l in enumerate(section_lessons) if l["id"] == lesson_id)
        pts = " | ".join(lesson["key_points"][:3])
        return f'{lead}\n\nLesson {num} - {lesson["title"]}:\n{pts}'

    def has(*keywords):
        return any(re.search(r"(?<!\w)" + re.escape(k) + r"(?!\w)", text) for k in keywords)

    def is_topic(*words):
        return any(re.search(r"(?<!\w)" + re.escape(w) + r"(?!\w)", text) for w in words)

    def code_example(message):
        def blocks(entries):
            return "\n\n".join(
                "```%s\n%s\n```" % (lang, code.rstrip()) for lang, code in entries
            )

        if is_topic("loop", "loops", "for loop", "while loop", "iteration"):
            return (
                "Here is a loop that prints the numbers 1 to 5 in all three languages — "
                "press Run on any block to execute it right here:\n\n"
                + blocks(
                    [
                        (
                            "python",
                            "for i in range(1, 6):\n    print(i)",
                        ),
                        (
                            "cpp",
                            "#include <iostream>\nusing namespace std;\n\n"
                            "int main() {\n"
                            "    for (int i = 1; i <= 5; i++) {\n"
                            "        cout << i << endl;\n"
                            "    }\n"
                            "    return 0;\n"
                            "}",
                        ),
                        (
                            "java",
                            "public class Main {\n"
                            "    public static void main(String[] args) {\n"
                            "        for (int i = 1; i <= 5; i++) {\n"
                            "            System.out.println(i);\n"
                            "        }\n"
                            "    }\n"
                            "}",
                        ),
                    ]
                )
            )
        if is_topic("function", "functions", "method", "def"):
            return (
                "Here is a function that adds two numbers and prints the result — "
                "press Run on any block:\n\n"
                + blocks(
                    [
                        (
                            "python",
                            "def add(a, b):\n"
                            "    return a + b\n\n"
                            "print(add(3, 4))",
                        ),
                        (
                            "cpp",
                            "#include <iostream>\nusing namespace std;\n\n"
                            "int add(int a, int b) {\n"
                            "    return a + b;\n"
                            "}\n\n"
                            "int main() {\n"
                            "    cout << add(3, 4) << endl;\n"
                            "    return 0;\n"
                            "}",
                        ),
                        (
                            "java",
                            "public class Main {\n"
                            "    static int add(int a, int b) {\n"
                            "        return a + b;\n"
                            "    }\n\n"
                            "    public static void main(String[] args) {\n"
                            "        System.out.println(add(3, 4));\n"
                            "    }\n"
                            "}",
                        ),
                    ]
                )
            )
        if is_topic("array", "arrays", "list", "lists", "vector", "collection"):
            return (
                "Here is a program that stores three colors in a collection and prints "
                "them — press Run on any block:\n\n"
                + blocks(
                    [
                        (
                            "python",
                            "colors = ['red', 'green', 'blue']\n"
                            "for c in colors:\n"
                            "    print(c)",
                        ),
                        (
                            "cpp",
                            "#include <iostream>\n#include <vector>\nusing namespace std;\n\n"
                            "int main() {\n"
                            "    vector<string> colors = {\"red\", \"green\", \"blue\"};\n"
                            "    for (string c : colors) {\n"
                            "        cout << c << endl;\n"
                            "    }\n"
                            "    return 0;\n"
                            "}",
                        ),
                        (
                            "java",
                            "import java.util.ArrayList;\n\n"
                            "public class Main {\n"
                            "    public static void main(String[] args) {\n"
                            "        ArrayList<String> colors = new ArrayList<>();\n"
                            "        colors.add(\"red\");\n"
                            "        colors.add(\"green\");\n"
                            "        colors.add(\"blue\");\n"
                            "        for (String c : colors) {\n"
                            "            System.out.println(c);\n"
                            "        }\n"
                            "    }\n"
                            "}",
                        ),
                    ]
                )
            )
        if is_topic("conditional", "if else", "if-else", "elif", "else if", "decision"):
            return (
                "Here is an if/else that checks if a number is even or odd — "
                "press Run on any block:\n\n"
                + blocks(
                    [
                        (
                            "python",
                            "n = 7\n"
                            "if n % 2 == 0:\n"
                            "    print('even')\n"
                            "else:\n"
                            "    print('odd')",
                        ),
                        (
                            "cpp",
                            "#include <iostream>\nusing namespace std;\n\n"
                            "int main() {\n"
                            "    int n = 7;\n"
                            "    if (n % 2 == 0) {\n"
                            "        cout << \"even\" << endl;\n"
                            "    } else {\n"
                            "        cout << \"odd\" << endl;\n"
                            "    }\n"
                            "    return 0;\n"
                            "}",
                        ),
                        (
                            "java",
                            "public class Main {\n"
                            "    public static void main(String[] args) {\n"
                            "        int n = 7;\n"
                            "        if (n % 2 == 0) {\n"
                            "            System.out.println(\"even\");\n"
                            "        } else {\n"
                            "            System.out.println(\"odd\");\n"
                            "        }\n"
                            "    }\n"
                            "}",
                        ),
                    ]
                )
            )
        if is_topic("variable", "variables", "data type", "datatype", "integer", "boolean", "float", "string", "char"):
            return (
                "Here is a program that stores a name and an age in variables and prints "
                "them — press Run on any block:\n\n"
                + blocks(
                    [
                        (
                            "python",
                            "name = 'Ana'\n"
                            "age = 20\n"
                            "print(name, age)",
                        ),
                        (
                            "cpp",
                            "#include <iostream>\nusing namespace std;\n\n"
                            "int main() {\n"
                            "    string name = \"Ana\";\n"
                            "    int age = 20;\n"
                            "    cout << name << \" \" << age << endl;\n"
                            "    return 0;\n"
                            "}",
                        ),
                        (
                            "java",
                            "public class Main {\n"
                            "    public static void main(String[] args) {\n"
                            "        String name = \"Ana\";\n"
                            "        int age = 20;\n"
                            "        System.out.println(name + \" \" + age);\n"
                            "    }\n"
                            "}",
                        ),
                    ]
                )
            )
        return (
            "Here is a starter program that prints a friendly message — "
            "press Run on any block to see it work:\n\n"
            + blocks(
                [
                    (
                        "python",
                        "print('Hello from CodeFundamentals!')",
                    ),
                    (
                        "cpp",
                        "#include <iostream>\nusing namespace std;\n\n"
                        "int main() {\n"
                        "    cout << \"Hello from CodeFundamentals!\" << endl;\n"
                        "    return 0;\n"
                        "}",
                    ),
                    (
                        "java",
                        "public class Main {\n"
                        "    public static void main(String[] args) {\n"
                        "        System.out.println(\"Hello from CodeFundamentals!\");\n"
                        "    }\n"
                        "}",
                    ),
                ]
            )
        )

    rules = [
        (lambda: has("long quiz", "section quiz"), "Each section ends with a Long Quiz: 15 questions with 20 seconds per question, and you get 2 attempts total. Your best score is saved on the dashboard. Unlock it by completing every lesson in the section, then find it on your dashboard under By Section, or use the Proceed to Long Quiz button after the last lesson. After taking it you can view the answer review to see which questions you got right or wrong."),
        (lambda: has("quiz", "quizzes", "attempt"), "Each lesson ends with a quiz. You get one attempt per quiz, so answer carefully! Your score is saved on your dashboard and added to the lesson and global leaderboards."),
        (lambda: has("leaderboard", "rank", "ranking", "points", "compete"), "Leaderboards rank users by total points: quiz scores plus lab bonuses. Completing a SQL Injection Lab earns 2 points per lab. There is a global leaderboard on your dashboard and a separate one for each lesson page."),
        (lambda: has("progress", "completed", "dashboard"), "Your dashboard shows your progress: lessons completed, quiz stats, and total points. Mark a lesson complete from the lesson page, then watch your progress bar grow."),
        (lambda: has("vs"), "You learn Python, C++, Java, and C side by side. Python is interpreted and uses indentation; C++ and Java are compiled and use braces {} and semicolons; C is a compiled language that is simpler and closer to the hardware. Every lesson shows the same concept in each language, so you can compare them directly in the lessons."),
        (lambda: has("c++", "cpp", "c plus plus"), "C++ is a compiled language: a compiler turns your source code into a fast machine program before it runs. It is used for games, operating systems, and performance-heavy software. C++ uses braces {} and requires semicolons at the end of each statement."),
        (lambda: has("java"), "Java is a compiled, object-oriented language that runs on the Java Virtual Machine (JVM), so the same code can run on many devices. It is widely used for Android apps and enterprise systems. Like C++, Java uses braces {} and semicolons, and you must declare variable types."),
        (lambda: has("python"), "Python is an interpreted language: it runs line by line, which makes it quick to test. Its readable syntax uses indentation instead of braces, which is why it is so beginner-friendly. Python is popular in data science, web apps, and automation."),
        (lambda: has("c language", "c programming", "c program", "learn c", "dennis ritchie", "printf"), "C is a compiled language created by Dennis Ritchie in 1972 at Bell Labs. It is the foundation of C++, Java, and Python, and is used for operating systems and embedded devices. Every C program needs a main() function and prints with printf(). C is taught side by side with Python, C++, and Java in every Fundamentals lesson, where you can run C code right in the browser."),
        (lambda: has("language", "languages", "syntax"), "You learn Python, C++, Java, and C side by side. Python is interpreted and uses indentation; C++ and Java are compiled and use braces {} and semicolons; C is a compact compiled language that inspired them both. Every lesson shows the same concept in each language."),
        (lambda: has("sql", "database", "dbms", "normaliz", "er model", "entity", "foreign key", "primary key", "ddl", "dml", "dql", "join"), "The Advanced Database System section covers databases and SQL: DBMS overview, ER model and keys, normalization, DDL (CREATE/ALTER/DROP), DML (INSERT/UPDATE/DELETE), SELECT and filtering, JOINs, and subqueries with aggregates. The SQL examples in those lessons run right in your browser against a live database engine."),
        (lambda: has("error", "debug", "compile", "bug"), "Debugging tips: read the error message first, check semicolons and braces in C++/Java, check indentation in Python, and add print statements to see what your code is doing. The code editor lets you test changes instantly."),
        (lambda: has("register", "sign up", "account", "free"), "Creating an account is free and takes seconds. It saves your progress, quiz scores, profile picture, and leaderboard standing."),
        (lambda: has("how it works", "how does this work", "site work"), "It is simple: create an account, read the lessons, run code live in the browser, take the quiz after each lesson, and climb the leaderboard. The Fundamentals section covers programming basics, the Advanced Fundamentals section teaches GitHub and deployment, and the Advanced Database System section covers databases and SQL."),
        (lambda: has("example code", "code example", "sample code", "show me code", "give me code", "write code", "code for", "snippet", "kodigo", "sample"), code_example(text)),
        (lambda: has("netlify", "netlify.com", "deploy netlify", "netlify deploy", "netlify.toml"), brief(33, "Netlify is a web platform for building and deploying modern web projects. Deploy static sites, SPAs, and JAMstack apps instantly from Git. Use a netlify.toml file to define build settings. The Deploying with Netlify lesson covers forms, functions, and custom domains.")),
        (lambda: has("render", "render.com", "deploy render", "render deploy"), brief(32, "Render is a cloud platform that deploys your app from GitHub in minutes. It handles SSL, auto-deploys, databases, and scaling. Use a render.yaml file to define your infrastructure. The Deploying with Render lesson covers the full setup.")),
        (lambda: has("github", "git", "repository", "repo", "commit", "branch", "merge", "pull request", "clone"), brief(31, "GitHub is the world's largest code hosting platform. Git tracks changes to your code, GitHub hosts repos online, and pull requests let you collaborate safely. The Introduction to GitHub lesson covers everything from git init to branching strategies.")),
        (lambda: has("deploy", "deployment", "deploying", "hosting", "production", "ci/cd", "ci cd", "vercel", "github pages"), brief(30, "Web Deployment Fundamentals covers deploying your app to the internet. Laravel needs a server with PHP and Nginx; React apps can be deployed to Vercel, Netlify, or Cloudflare Pages. Always use environment variables for secrets and set APP_DEBUG=false in production.")),
        (lambda: has("variable", "variables", "data type", "datatype", "integer", "boolean", "float", "string", "char"), brief(2, "Variables store values in memory. Python infers the type automatically; C++ and Java require you to declare it before the name.")),
        (lambda: has("logical"), brief(5, "Logical operators combine true/false values. Python: and, or, not. C++/Java: &&, ||, !. Both conditions are checked to build bigger conditions.")),
        (lambda: has("operator", "operators", "arithmetic", "modulo", "remainder"), brief(4, "Operators do math and comparisons. Arithmetic: + - * / %. Relational: == != < > <= >=. Comparison always produces true or false.")),
        (lambda: has("conditional", "if else", "if-else", "elif", "else if", "decision"), brief(6, "Conditionals make decisions. Python uses if/elif/else with colons; C++ and Java use if/else if/else with braces and parentheses.")),
        (lambda: has("loop", "loops", "for loop", "while", "repeat", "iteration"), brief(7, "Loops repeat a block of code. For loops run a fixed number of times; while loops repeat while a condition is true. Beware of infinite loops!")),
        (lambda: has("function", "functions", "methods", "def", "return"), brief(8, "A function is a reusable block of code. It takes parameters and can return a result. Python uses def; C++/Java declare a return type.")),
        (lambda: has("array", "arrays", "list", "lists", "vector", "collection", "index"), brief(9, "Collections store multiple values. Python uses lists, C++ arrays/vectors, Java arrays/ArrayLists. Indexing always starts at 0.")),
        (lambda: has("problem solving", "solve problem", "pseudocode", "plan"), brief(10, "Problem solving is a process: understand the problem, plan the steps (pseudocode), write the code, then test with simple examples.")),
        (lambda: has("programming", "what is code", "what is a program"), brief(1, "Programming is writing instructions that a computer can follow. Source code is human-readable, then translated or interpreted for the machine.")),
        (lambda: has("who are you", "your name", "eren", "story", "about you"), "I am E.R.E.N — Educational Response Engine for Novices. My story started at 2 AM in a dorm room: a frustrated first-year CS student kept failing a quiz about loops, so they built me to explain concepts the way a patient friend would. Every question a classmate asked that I could not answer became a new rule, and here I am. Ask me anything!"),
        (lambda: has("hi", "hello", "hey", "good morning", "good afternoon", "good evening"), "Hello! I am E.R.E.N, your CodeFundamentals assistant. Ask me about variables, operators, logical operators, conditionals, loops, functions, arrays, the long quiz, quizzes, or the leaderboard."),
        (lambda: has("thanks", "thank you", "salamat"), "You are welcome! Happy coding."),
        (lambda: has("help", "what can you do"), "I can explain lessons and concepts, quiz and long quiz rules, leaderboard rules, debugging tips, and how the site works. Try asking: 'explain loops', 'how does the long quiz work?', or 'how does the quiz work?'"),
    ]
    for condition, reply in rules:
        if condition():
            return reply
    knowledge_reply = general_knowledge_reply(text)
    if knowledge_reply:
        return knowledge_reply
    if is_definition_question(text):
        return (
            "Good question, but that one is outside my built-in course notes so I cannot "
            "explain it well yet. Ask me about variables, operators, logical operators, "
            "conditionals, loops, functions, arrays, quizzes, the long quiz, or the "
            "leaderboard — or when my online AI brain is connected I can answer general "
            "questions like 'what is research?'."
        )
    return (
        "I am not sure about that one yet. Try asking about variables, operators, logical "
        "operators, conditionals, loops, functions, arrays, quizzes, the long quiz, or the "
        "leaderboard — and if it is a general question like 'what is research?', I will do my "
        "best to answer it."
    )


def extract_youtube_id(text):
    m = re.search(
        r"(?:https?://)?(?:www\.)?(?:youtube\.com/(?:watch\?v=|embed/|shorts/|live/|v/)|youtu\.be/)([A-Za-z0-9_-]{11})",
        text or "",
    )
    return m.group(1) if m else None


def fetch_youtube_transcript(video_id, max_chars=120000):
    try:
        from youtube_transcript_api import YouTubeTranscriptApi

        def _part_text(part):
            if isinstance(part, dict):
                return part.get("text", "")
            return getattr(part, "text", "") or ""

        try:
            api = YouTubeTranscriptApi()
            fetched = api.fetch(
                video_id,
                languages=["en", "en-US", "en-GB", "fil", "tl", "es"],
            )
        except Exception:
            try:
                api = YouTubeTranscriptApi()
                fetched = api.fetch(video_id)
            except AttributeError:
                fetched = YouTubeTranscriptApi.get_transcript(video_id)
        text = " ".join(_part_text(p) for p in fetched)
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            return None
        return text[:max_chars]
    except Exception as exc:
        print(
            "[E.R.E.N] Transcript fetch failed for %s: %s" % (video_id, exc),
            file=sys.stderr,
        )
        return None


def gemini_summarize(prompt):
    api_key = _gemini_api_key()
    if not api_key:
        return None
    model = _gemini_model()
    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        f"{model}:generateContent?key={api_key}"
    )
    payload = {
        "system_instruction": {
            "parts": [{"text": EREN_SYSTEM_PROMPT}]
        },
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.4, "maxOutputTokens": 4096},
    }
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            body = json.loads(resp.read().decode("utf-8"))
        parts = body.get("candidates", [{}])[0].get("content", {}).get("parts") or []
        reply = "".join(p.get("text", "") for p in parts).strip()
        return reply or None
    except Exception as exc:
        print(
            f"[E.R.E.N] Gemini summarize request failed ({model}): {exc}",
            file=sys.stderr,
        )
        return None


def youtube_summary_reply(video_id, user_message):
    transcript = fetch_youtube_transcript(video_id)
    if not transcript:
        return (
            "Hindi ko makuha ang transcript ng video na iyan — maaaring naka-disable ang "
            "captions nito o hindi available ang video. Subukan ang ibang video, o i-paste "
            "ang script sa chat para i-summarize ko."
        )
    prompt = (
        "A student pasted a YouTube link (video ID: %s) and wants help with it. Their "
        "message: \"%s\"\n\n"
        "Here is the video's transcript:\n\n%s\n\n"
        "Give a clear, beginner-friendly summary of the video in the student's language "
        "(English or Filipino). Use short sections with headers and bullet points "
        "(## and - only, no other markdown), highlight the key points, and answer their "
        "question if they asked one. Keep it under 700 words."
    ) % (video_id, user_message[:300], transcript)
    reply = gemini_summarize(prompt)
    if reply:
        return reply
    preview = transcript[:800]
    return (
        "Nakuha ko ang transcript ng video, pero walang available na AI brain (missing "
        "GEMINI_API_KEY), kaya ito lang ang preview ng mga unang bahagi:\n\n"
        "> %s...\n\n"
        "I-set ang GEMINI_API_KEY para makapagbigay ako ng buong buod." % preview
    )


@app.route("/assistant", methods=["POST"])
@login_required
def assistant():
    data = request.get_json(silent=True) or {}
    message = (data.get("message") or "").strip()
    if not message:
        return jsonify({"ok": False, "error": "Please type a message."}), 400
    video_id = extract_youtube_id(message)
    if video_id:
        return jsonify({"ok": True, "reply": youtube_summary_reply(video_id, message), "yt_summary": True})
    history = data.get("history") if isinstance(data.get("history"), list) else None
    reply = gemini_reply(message, history)
    if reply is None:
        reply = assistant_reply(message)
    return jsonify({"ok": True, "reply": reply})


def clean_inline_text(s):
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    s = re.sub(r"^\*+\s*", "", s)
    s = re.sub(r"^Title\s*:\s*", "", s.strip(), flags=re.IGNORECASE)
    s = re.sub(r"^Slide\s*\d+\s*[:\-.\u2013]?\s*", "", s.strip(), flags=re.IGNORECASE)
    return s.strip()


def parse_slides_text(text):
    slides = []
    current = None
    in_code = False
    code_buf = []
    code_lang = ""

    def ensure_slide(title):
        nonlocal current
        current = {"title": clean_inline_text(title) or "Presentation", "items": []}
        slides.append(current)

    for raw in text.splitlines():
        line = raw.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                if current is not None:
                    code_text = "\n".join(code_buf).strip()
                    if code_lang == "output":
                        current["items"].append({"type": "output", "text": code_text})
                    else:
                        current["items"].append({"type": "code", "text": code_text})
                code_buf = []
                in_code = False
                code_lang = ""
            else:
                in_code = True
                code_buf = []
                code_lang = stripped[3:].strip().lower()
            continue

        if in_code:
            code_buf.append(raw)
            continue

        if not stripped:
            continue

        if re.match(r"^\*{0,2}(sample\s+)?output\s*:\s*\*{0,2}$", stripped, re.IGNORECASE):
            continue

        hm = re.match(r"^#{1,6}\s+(.*)$", stripped)
        if hm:
            ensure_slide(hm.group(1))
            continue

        sm = re.match(r"^Slide\s*\d+\s*([:\-.\u2013]\s*)?(.*)$", stripped, re.IGNORECASE)
        if sm and sm.group(2):
            ensure_slide(sm.group(2))
            continue

        if current is None:
            tm = re.match(r"^Title\s*:\s*(.*)$", stripped, re.IGNORECASE)
            ensure_slide(tm.group(1) if tm else "Overview")
            if tm:
                continue

        bm = re.match(r"^([\-\*\+])\s+(.*)$", stripped)
        if bm:
            text = bm.group(2)
            if re.match(r"^(Bullets?|(sample\s+)?output)\s*:\s*$", text, re.IGNORECASE):
                continue
            indent = len(line) - len(line.lstrip())
            current["items"].append(
                {"text": clean_inline_text(text), "level": min(indent // 2, 3)}
            )
            continue

        nm = re.match(r"^(\d+)[\.\)]\s+(.*)$", stripped)
        if nm:
            current["items"].append({"text": clean_inline_text(nm.group(2)), "level": 0})
            continue

        current["items"].append({"text": clean_inline_text(stripped), "level": 0})

    if not slides:
        slides = [{"title": "Presentation", "items": []}]
    return slides


def build_pptx_export(text):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    slides = parse_slides_text(text)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    NAVY = RGBColor(15, 23, 42)
    SLATE = RGBColor(51, 65, 85)
    BRAND = RGBColor(14, 165, 233)
    CODE_BG = RGBColor(15, 23, 42)
    CODE_FG = RGBColor(148, 210, 255)
    CODE_NAME = "Consolas"
    OUT_BG = RGBColor(6, 46, 30)
    OUT_LINE = RGBColor(16, 185, 129)
    OUT_FG = RGBColor(167, 243, 208)
    OUT_HEAD = RGBColor(52, 211, 153)
    FONT = "Calibri"

    for slide in slides:
        s = prs.slides.add_slide(layout)
        bullets = [it for it in slide["items"] if it.get("type") not in ("code", "output")]
        codes = [it for it in slide["items"] if it.get("type") == "code"]
        outputs = [it for it in slide["items"] if it.get("type") == "output"]

        title = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
        tf = title.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = 0
        tf.margin_top = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = slide["title"]
        r.font.name = FONT
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = NAVY

        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.52), Inches(2.2), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BRAND
        bar.line.fill.background()

        content_top = Inches(1.95)
        if codes:
            content_h = Inches(2.5)
        else:
            content_h = Inches(5.2)
        if bullets:
            box = s.shapes.add_textbox(Inches(0.65), content_top, Inches(12.05), content_h)
            btf = box.text_frame
            btf.word_wrap = True
            btf.vertical_anchor = MSO_ANCHOR.TOP
            btf.margin_left = 0
            btf.margin_top = 0
            for i, it in enumerate(bullets):
                bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                bp.alignment = PP_ALIGN.LEFT
                marker = "\u2022  " if it["level"] == 0 else "\u2013  "
                indent = "     " * min(it["level"], 2) if it["level"] > 0 else ""
                bp.text = marker + indent + it["text"]
                for br in bp.runs:
                    br.font.name = FONT
                    br.font.size = Pt(max(17 - it["level"] * 2, 14))
                    br.font.color.rgb = SLATE

        if codes:
            top = Inches(4.55)
            bottom = Inches(7.3)
            gap = Inches(0.14)
            units = []
            total = 0
            for i, code in enumerate(codes):
                out_text = (outputs[i].get("text") if i < len(outputs) else "") or ""
                code_lines = code["text"].split("\n")
                out_lines = out_text.split("\n") if out_text.strip() else []
                cl = len(code_lines)
                ol = len(out_lines)
                csize = max(10, min(13, int(210 / max(cl, 1))))
                osize = max(9, min(12, int(160 / max(ol, 1)))) if ol else 11
                ch = int(Inches(0.24)) + int(Pt(csize * 1.45) * cl)
                oh = (int(Inches(0.3)) + int(Pt(osize * 1.45) * (ol + 1))) if ol else 0
                units.append({
                    "code_lines": code_lines, "out_lines": out_lines,
                    "cl": cl, "ol": ol, "csize": csize, "osize": osize,
                    "ch": ch, "oh": oh,
                })
                total += ch + oh + int(gap)
            avail = int(bottom - top)
            scale = 1.0
            if total > avail:
                scale = (avail - int(gap)) / total if total else 1.0
            for u in units:
                ch = Emu(int(u["ch"] * scale))
                oh = Emu(int(u["oh"] * scale))
                csize = max(8, int(u["csize"] * scale))
                osize = max(8, int(u["osize"] * scale))
                box = s.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top, Inches(12.05), ch
                )
                try:
                    box.adjustments[0] = 0.06
                except Exception:
                    pass
                box.fill.solid()
                box.fill.fore_color.rgb = CODE_BG
                box.line.fill.background()
                box.shadow.inherit = False
                ctf = box.text_frame
                ctf.word_wrap = True
                ctf.vertical_anchor = MSO_ANCHOR.TOP
                ctf.margin_left = Inches(0.25)
                ctf.margin_right = Inches(0.25)
                ctf.margin_top = Inches(0.1)
                ctf.margin_bottom = Inches(0.06)
                for i, ln in enumerate(u["code_lines"]):
                    cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                    cp.alignment = PP_ALIGN.LEFT
                    cp.text = ln or " "
                    for cr in cp.runs:
                        cr.font.name = CODE_NAME
                        cr.font.size = Pt(csize)
                        cr.font.color.rgb = CODE_FG
                if u["ol"]:
                    oy = Emu(int(top) + int(ch) + int(gap))
                    obox = s.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), oy, Inches(12.05), oh
                    )
                    try:
                        obox.adjustments[0] = 0.06
                    except Exception:
                        pass
                    obox.fill.solid()
                    obox.fill.fore_color.rgb = OUT_BG
                    obox.line.color.rgb = OUT_LINE
                    obox.line.width = Pt(1)
                    obox.shadow.inherit = False
                    otf = obox.text_frame
                    otf.word_wrap = True
                    otf.vertical_anchor = MSO_ANCHOR.TOP
                    otf.margin_left = Inches(0.25)
                    otf.margin_right = Inches(0.25)
                    otf.margin_top = Inches(0.1)
                    otf.margin_bottom = Inches(0.06)
                    hdr = otf.paragraphs[0]
                    hdr.alignment = PP_ALIGN.LEFT
                    hr = hdr.add_run()
                    hr.text = "Sample Output"
                    hr.font.name = FONT
                    hr.font.size = Pt(11)
                    hr.font.bold = True
                    hr.font.color.rgb = OUT_HEAD
                    for i, ln in enumerate(u["out_lines"]):
                        cp = otf.add_paragraph()
                        cp.alignment = PP_ALIGN.LEFT
                        cp.text = ln or " "
                        for cr in cp.runs:
                            cr.font.name = CODE_NAME
                            cr.font.size = Pt(osize)
                            cr.font.color.rgb = OUT_FG
                    top = Emu(int(oy) + int(oh))
                else:
                    top = Emu(int(top) + int(ch) + int(gap))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route("/assistant/ppt", methods=["POST"])
@login_required
def assistant_ppt():
    data = request.get_json(silent=True) or {}
    text = (data.get("text") or "").strip()
    if not text:
        return jsonify({"ok": False, "error": "Nothing to export."}), 400
    try:
        slides = parse_slides_text(text)
        topic = (slides[0] or {}).get("title") or "presentation"
        base = re.sub(r"[^\w \-]+", "", topic).strip().replace(" ", "-") or "presentation"
        buf = build_pptx_export(text)
    except Exception as exc:
        return jsonify({"ok": False, "error": "Could not create the file: %s" % exc}), 500
    return send_file(
        buf,
        as_attachment=True,
        download_name=(base[:60] + ".pptx"),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@app.route("/lesson/<int:lesson_id>")
@login_required
def lesson(lesson_id):
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None:
        flash("Lesson not found.", "danger")
        return redirect(url_for("dashboard"))
    if lesson.get("coming_soon"):
        flash("SQL Injection is coming soon — check back later!", "info")
        return redirect(url_for("dashboard"))
    section_lessons = [
        l
        for l in LESSONS
        if l.get("section", "Fundamentals") == lesson.get("section", "Fundamentals")
        and not l.get("coming_soon")
    ]
    idx = next(i for i, l in enumerate(section_lessons) if l["id"] == lesson_id)
    completed_ids = get_completed_ids(current_user)
    unlocked_ids = get_unlocked_ids(current_user)
    if (
        idx > 0
        and section_lessons[idx - 1]["id"] not in completed_ids
        and lesson_id not in unlocked_ids
    ):
        flash(
            f'Complete Lesson {idx} ({section_lessons[idx - 1]["title"]}) first — or unlock this lesson with coins — to access it.',
            "danger",
        )
        return redirect(url_for("lesson", lesson_id=section_lessons[idx - 1]["id"]))
    completed = lesson_id in completed_ids
    quiz_score = get_quiz_scores(current_user).get(lesson_id)
    has_lab = lesson_id in LAB_LESSON_IDS
    labs_leaderboard = get_lesson_labs_leaderboard(lesson_id) if has_lab else []
    lab_completed = is_lab_completed(current_user, lesson_id) if has_lab else False
    lab_url = None
    lab_title = None
    if has_lab:
        if lesson_id in CUSTOM_LAB_ROUTES:
            lab_url = url_for(CUSTOM_LAB_ROUTES[lesson_id])
        else:
            lab_url = url_for("lesson_lab", lesson_id=lesson_id)
        lab_title = LABS[lesson_id]["title"]
    prev_lesson_id = section_lessons[idx - 1]["id"] if idx > 0 else None
    next_lesson_id = section_lessons[idx + 1]["id"] if idx < len(section_lessons) - 1 else None
    return render_template(
        "lesson.html",
        lesson=lesson,
        est_minutes=lesson_estimate_minutes(lesson),
        lesson_num=idx + 1,
        section_total=len(section_lessons),
        completed=completed,
        quiz_score=quiz_score,
        lesson_leaderboard=get_lesson_leaderboard(lesson_id),
        has_lab=has_lab,
        labs_leaderboard=labs_leaderboard,
        lab_completed=lab_completed,
        lab_url=lab_url,
        lab_title=lab_title,
        total_lessons=len(AVAILABLE_LESSONS),
        prev_lesson_id=prev_lesson_id,
        next_lesson_id=next_lesson_id,
        section_minutes=SECTION_LESSON_MINUTES,
        section_has_quiz=lesson.get("section", "Fundamentals") in LONG_QUIZZES,
        section_lessons=section_lessons,
        completed_ids=completed_ids,
        unlocked_ids=unlocked_ids,
    )


@app.route("/lesson/<int:lesson_id>/leaderboard", methods=["GET"])
@login_required
def lesson_leaderboard(lesson_id):
    if not any(l["id"] == lesson_id for l in LESSONS):
        return jsonify({"ok": False, "error": "Lesson not found."}), 404
    return jsonify({"ok": True, "leaderboard": get_lesson_leaderboard(lesson_id)})


def _activity_docx_path(lesson):
    docx_name = (lesson.get("activity") or {}).get("docx")
    if not docx_name:
        return None
    path = os.path.join(app.static_folder, "docs", docx_name)
    if not os.path.exists(path):
        return None
    return path


@app.route("/lesson/<int:lesson_id>/activity/questions", methods=["GET"])
@login_required
def activity_questions(lesson_id):
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None or not lesson.get("activity"):
        return jsonify({"ok": False, "error": "No activity for this lesson."}), 404
    path = _activity_docx_path(lesson)
    if path is None:
        return jsonify({"ok": False, "error": "Activity file is missing."}), 404
    from docx import Document as DocxDocument

    doc = DocxDocument(path)
    in_questions = False
    questions = []
    meta = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text == "Questions":
            in_questions = True
            continue
        if in_questions and text:
            questions.append(text)
    for t in doc.tables:
        for row in t.rows:
            cells = row.cells
            if len(cells) >= 2 and cells[0].text.strip():
                label = cells[0].text.strip().rstrip(":").strip()
                meta.append({"label": label, "key": label})
    return jsonify({
        "ok": True,
        "title": lesson["activity"].get("title", "Activity"),
        "meta": meta,
        "questions": questions,
    })


@app.route("/lesson/<int:lesson_id>/activity/download", methods=["POST"])
@login_required
def activity_download(lesson_id):
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None or not lesson.get("activity"):
        return jsonify({"ok": False, "error": "No activity for this lesson."}), 404
    path = _activity_docx_path(lesson)
    if path is None:
        return jsonify({"ok": False, "error": "Activity file is missing."}), 404
    data = request.get_json(silent=True) or {}
    answers = data.get("answers") or {}
    meta_values = data.get("meta") or {}
    if not isinstance(answers, dict):
        answers = {}
    if not isinstance(meta_values, dict):
        meta_values = {}

    from copy import deepcopy
    from docx import Document as DocxDocument
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    def insert_answer_after(paragraph, text, indent_twips=720):
        anchor = paragraph._p
        pPr_source = None
        nxt = anchor.getnext()
        if nxt is not None:
            nxt_pPr = nxt.find(qn("w:pPr"))
            if nxt_pPr is not None and nxt_pPr.find(qn("w:numPr")) is None:
                pPr_source = nxt_pPr
        mark_fonts = None
        mark_pPr = anchor.find(qn("w:pPr"))
        if mark_pPr is not None:
            mr = mark_pPr.find(qn("w:rPr"))
            if mr is not None:
                mark_fonts = mr.find(qn("w:rFonts"))
        for line in text.splitlines():
            new_p = deepcopy(anchor)
            for child in list(new_p):
                new_p.remove(child)
            pPr = deepcopy(pPr_source) if pPr_source is not None else OxmlElement("w:pPr")
            numPr = pPr.find(qn("w:numPr"))
            if numPr is not None:
                pPr.remove(numPr)
            ind = OxmlElement("w:ind")
            ind.set(qn("w:left"), str(indent_twips))
            rPr_node = pPr.find(qn("w:rPr"))
            if rPr_node is not None:
                rPr_node.addprevious(ind)
            else:
                pPr.append(ind)
            new_p.append(pPr)
            r = OxmlElement("w:r")
            rPr = OxmlElement("w:rPr")
            for tag in ("w:b", "w:i", "w:u"):
                el = OxmlElement(tag)
                el.set(qn("w:val"), "0")
                rPr.append(el)
            if mark_fonts is not None:
                rPr.append(deepcopy(mark_fonts))
            r.append(rPr)
            t = OxmlElement("w:t")
            t.text = line if line else " "
            t.set(qn("xml:space"), "preserve")
            r.append(t)
            new_p.append(r)
            pPr_source = pPr
            anchor.addnext(new_p)
            anchor = new_p

    doc = DocxDocument(path)
    in_questions = False
    qi = 0
    for p in doc.paragraphs:
        text = p.text.strip()
        if text == "Questions":
            in_questions = True
            continue
        if in_questions and text:
            answer = answers.get(str(qi)) or ""
            answer = answer.strip()
            if answer:
                insert_answer_after(p, answer)
            qi += 1

    for t in doc.tables:
        for row in t.rows:
            cells = row.cells
            if len(cells) >= 2:
                label = cells[0].text.strip().rstrip(":").strip()
                value = meta_values.get(label) or ""
                if label and value.strip():
                    p = cells[1].paragraphs[0]
                    if p.runs:
                        p.runs[0].text = value.strip()
                    else:
                        p.add_run(value.strip())

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    base = os.path.splitext(os.path.basename(path))[0]
    safe_name = re.sub(r"[^A-Za-z0-9_. -]", "", current_user.username).strip() or "user"
    return send_file(
        buf,
        as_attachment=True,
        download_name="%s - %s.docx" % (base, safe_name),
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@app.route("/lesson/<int:lesson_id>/quiz", methods=["POST"])
@login_required
def submit_quiz(lesson_id):
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None:
        return jsonify({"ok": False, "error": "Lesson not found."}), 404
    data = request.get_json(silent=True) or {}
    try:
        score = int(data.get("score"))
        total = int(data.get("total"))
    except (TypeError, ValueError):
        return jsonify({"ok": False, "error": "Invalid score."}), 400
    expected_total = len(lesson["quiz"])
    if total != expected_total or not (0 <= score <= total):
        return jsonify({"ok": False, "error": "Invalid score."}), 400

    row = QuizScore.query.filter_by(
        user_id=current_user.id, lesson_id=lesson_id
    ).first()
    if row:
        return jsonify(
            {"ok": False, "error": "You already took this quiz (1 attempt only)."}
        ), 400
    row = QuizScore(
        user_id=current_user.id,
        lesson_id=lesson_id,
        best_score=score,
        total_questions=total,
        attempts=1,
    )
    db.session.add(row)
    db.session.commit()
    log_activity(
        current_user,
        "Completed lesson quiz",
        "Scored %d/%d on \"%s\"" % (score, total, lesson["title"]),
    )
    return jsonify(
        {"ok": True, "best_score": row.best_score, "attempts": row.attempts}
    )


@app.route("/lesson/<int:lesson_id>/unlock", methods=["POST"])
@login_required
def unlock_lesson(lesson_id):
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None or lesson.get("coming_soon"):
        flash("Lesson not found.", "danger")
        return redirect(url_for("dashboard"))
    unlocked_ids = get_unlocked_ids(current_user)
    if lesson_id in unlocked_ids:
        flash("You already unlocked this lesson.", "info")
        return redirect(url_for("dashboard"))
    cost = lesson_unlock_cost(lesson_id)
    if (current_user.coins or 0) < cost:
        flash(
            "You need %d coins to unlock this lesson. Earn coins by taking quizzes and completing lessons."
            % cost,
            "warning",
        )
        return redirect(url_for("dashboard"))
    current_user.coins = (current_user.coins or 0) - cost
    db.session.add(UnlockedLesson(user_id=current_user.id, lesson_id=lesson_id))
    db.session.commit()
    log_activity(
        current_user,
        "Unlocked lesson with coins",
        'Unlocked "%s" (-%d coins)' % (lesson["title"], cost),
    )
    flash(
        "Lesson unlocked! You spent %d coins." % cost,
        "success",
    )
    return redirect(url_for("lesson", lesson_id=lesson_id))


@app.route("/lesson/<int:lesson_id>/complete", methods=["POST"])
@login_required
def complete_lesson(lesson_id):
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None:
        flash("Lesson not found.", "danger")
        return redirect(url_for("dashboard"))
    quiz_row = QuizScore.query.filter_by(
        user_id=current_user.id, lesson_id=lesson_id
    ).first()
    if quiz_row is None:
        flash(
            "Take the lesson quiz first before marking it as complete.",
            "warning",
        )
        return redirect(url_for("lesson", lesson_id=lesson_id))
    exists = LessonProgress.query.filter_by(
        user_id=current_user.id, lesson_id=lesson_id
    ).first()
    if not exists:
        db.session.add(LessonProgress(user_id=current_user.id, lesson_id=lesson_id))
        current_user.coins = (current_user.coins or 0) + COINS_PER_LESSON
        db.session.commit()
        log_activity(
            current_user,
            "Completed lesson",
            'Finished "%s" (+%d coins)' % (lesson["title"], COINS_PER_LESSON),
        )
        flash(
            "Lesson marked as complete. +%d coins earned!" % COINS_PER_LESSON,
            "success",
        )
    else:
        flash("You already completed this lesson.", "info")
    return redirect(url_for("lesson", lesson_id=lesson_id))


@app.route("/sql-lab")
@login_required
def sql_lab():
    return render_template(
        "sql_lab.html",
        lab_completed=is_lab_completed(current_user, 19),
        lesson_id=19,
    )


@app.route("/sql-lab/complete", methods=["POST"])
@login_required
def sql_lab_complete():
    mark_lab_completed(current_user, 19)
    log_activity(current_user, "Completed lab", "Finished SQL Injection Lab 1")
    return jsonify({"ok": True})


@app.route("/sql-lab/reset", methods=["POST"])
@login_required
def sql_lab_reset():
    unmark_lab_completed(current_user, 19)
    log_activity(current_user, "Reset lab", "Reset SQL Injection Lab 1")
    return jsonify({"ok": True})


@app.route("/sql-lab-2")
@login_required
def sql_lab_2():
    return render_template(
        "sql_lab_2.html",
        lab_completed=is_lab_completed(current_user, 20),
        lesson_id=20,
    )


@app.route("/sql-lab-2/complete", methods=["POST"])
@login_required
def sql_lab_2_complete():
    mark_lab_completed(current_user, 20)
    log_activity(current_user, "Completed lab", "Finished SQL Injection Lab 2")
    return jsonify({"ok": True})


@app.route("/sql-lab-2/reset", methods=["POST"])
@login_required
def sql_lab_2_reset():
    unmark_lab_completed(current_user, 20)
    log_activity(current_user, "Reset lab", "Reset SQL Injection Lab 2")
    return jsonify({"ok": True})


@app.route("/lab/<int:lesson_id>")
@login_required
def lesson_lab(lesson_id):
    if lesson_id in CUSTOM_LAB_ROUTES:
        return redirect(url_for(CUSTOM_LAB_ROUTES[lesson_id]))
    lab = LABS.get(lesson_id)
    if lab is None:
        flash("Lab not found.", "danger")
        return redirect(url_for("dashboard"))
    lesson = next((l for l in LESSONS if l["id"] == lesson_id), None)
    if lesson is None:
        flash("Lesson not found.", "danger")
        return redirect(url_for("dashboard"))
    lab_json = json.dumps(lab).replace("</", "<\\/")
    return render_template(
        "lab.html",
        lab=lab,
        lab_json=lab_json,
        lesson=lesson,
        lesson_id=lesson_id,
        lab_completed=is_lab_completed(current_user, lesson_id),
    )


@app.route("/lab/<int:lesson_id>/complete", methods=["POST"])
@login_required
def lesson_lab_complete(lesson_id):
    lab = LABS.get(lesson_id)
    if lab is None:
        return jsonify({"ok": False, "error": "Lab not found."}), 404
    mark_lab_completed(current_user, lesson_id)
    log_activity(current_user, "Completed lab", "Finished \"%s\"" % lab["title"])
    return jsonify({"ok": True})


@app.route("/lab/<int:lesson_id>/reset", methods=["POST"])
@login_required
def lesson_lab_reset(lesson_id):
    lab = LABS.get(lesson_id)
    if lab is None:
        return jsonify({"ok": False, "error": "Lab not found."}), 404
    unmark_lab_completed(current_user, lesson_id)
    log_activity(current_user, "Reset lab", "Reset \"%s\"" % lab["title"])
    return jsonify({"ok": True})


@app.route("/section-quiz/<section>")
@login_required
def section_quiz(section):
    questions = LONG_QUIZZES.get(section)
    if not questions:
        flash("Long quiz not found.", "danger")
        return redirect(url_for("dashboard"))
    score = SectionQuizScore.query.filter_by(
        user_id=current_user.id, section=section
    ).first()
    last_answers = None
    if score and score.last_answers:
        try:
            last_answers = json.loads(score.last_answers)
        except (ValueError, TypeError):
            last_answers = None
    return render_template(
        "section_quiz.html",
        section=section,
        questions=questions,
        score=score,
        last_answers=last_answers,
    )


@app.route("/section-quiz/<section>/submit", methods=["POST"])
@login_required
def submit_section_quiz(section):
    questions = LONG_QUIZZES.get(section)
    if not questions:
        return jsonify({"ok": False, "error": "Long quiz not found."}), 404
    data = request.get_json(silent=True) or {}
    answers = data.get("answers")
    if not isinstance(answers, list) or len(answers) != len(questions):
        return jsonify({"ok": False, "error": "Invalid answers."}), 400
    score = 0
    for chosen, q in zip(answers, questions):
        if isinstance(chosen, int) and chosen == q["correct"]:
            score += 1
    row = SectionQuizScore.query.filter_by(
        user_id=current_user.id, section=section
    ).first()
    if row and row.attempts >= 2:
        return jsonify(
            {"ok": False, "error": "You already used your 2 attempts for this long quiz."}
        ), 400
    if row:
        row.best_score = max(row.best_score, score)
        row.attempts += 1
        row.last_answers = json.dumps(answers)
        db.session.commit()
    else:
        row = SectionQuizScore(
            user_id=current_user.id,
            section=section,
            best_score=score,
            total_questions=len(questions),
            attempts=1,
            last_answers=json.dumps(answers),
        )
        db.session.add(row)
        db.session.commit()
    log_activity(
        current_user,
        "Completed long quiz",
        "Scored %d/%d on the %s quiz" % (score, len(questions), section),
    )
    return jsonify(
        {
            "ok": True,
            "score": score,
            "best_score": row.best_score,
            "total_questions": row.total_questions,
            "attempts": row.attempts,
        }
    )


@app.route("/section-quiz/<section>/rate", methods=["POST"])
@login_required
def rate_section_quiz(section):
    questions = LONG_QUIZZES.get(section)
    if not questions:
        return jsonify({"ok": False, "error": "Long quiz not found."}), 404
    row = SectionQuizScore.query.filter_by(
        user_id=current_user.id, section=section
    ).first()
    if not row:
        return jsonify({"ok": False, "error": "Complete the long quiz first."}), 400
    data = request.get_json(silent=True) or {}
    ratings = data.get("ratings")
    if not isinstance(ratings, list) or len(ratings) != len(questions):
        return jsonify({"ok": False, "error": "Invalid ratings."}), 400
    cleaned = []
    for r in ratings:
        if not isinstance(r, dict):
            return jsonify({"ok": False, "error": "Invalid ratings."}), 400
        rating = r.get("rating")
        if not isinstance(rating, int) or rating < 1 or rating > 5:
            return jsonify({"ok": False, "error": "Rate every topic between 1 and 5."}), 400
        cleaned.append({"topic": str(r.get("topic", "")), "rating": rating})
    row.topic_ratings = json.dumps(cleaned)
    db.session.commit()
    log_activity(
        current_user,
        "Rated long quiz topics",
        "Submitted topic difficulty ratings for the %s quiz" % section,
    )
    return jsonify({"ok": True, "count": len(cleaned)})


RUN_CONFIGS = {
    "python": {
        "source": "main.py",
    },
    "cpp": {
        "source": "main.cpp",
    },
    "c": {
        "source": "main.c",
    },
    "java": {
        "source": "Main.java",
    },
    "typescript": {
        "source": "main.ts",
    },
}

COMPILER_FALLBACKS = {
    "g++": [r"C:\Users\chris\mingw64\bin\g++.exe"],
    "gcc": [r"C:\Users\chris\mingw64\bin\gcc.exe"],
    "clang++": [r"C:\Program Files\LLVM\bin\clang++.exe"],
    "javac": [r"C:\Users\chris\jdk21\jdk-21.0.12+8\bin\javac.exe"],
    "java": [r"C:\Users\chris\jdk21\jdk-21.0.12+8\bin\java.exe"],
}

MINGW_BIN = r"C:\Users\chris\mingw64\bin"

# C++/C online fallback: Compiler Explorer (free, keyless). Optional Judge0 backup
# via CODE_API_KEY env var: https://rapidapi.com/judge0-official/api/judge0-ce
CE_URL = "https://godbolt.org/api/compiler/g122/compile"
C_CE_URL = "https://godbolt.org/api/compiler/cg122/compile"
JAVA_CE_URL = "https://godbolt.org/api/compiler/java2102/compile"
JUDGE0_URL = "https://judge0-ce.p.rapidapi.com/submissions?base64_encoded=false&wait=true"


def find_tool(name):
    found = shutil.which(name)
    if found:
        return found
    for path in COMPILER_FALLBACKS.get(name, []):
        if os.path.isfile(path):
            return path
    return None


def run_locally(cmd, workdir, timeout, extra_env=None, stdin=""):
    env = dict(os.environ)
    if extra_env:
        env["PATH"] = extra_env + os.pathsep + env.get("PATH", "")
    flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
    return subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=timeout,
        cwd=workdir,
        creationflags=flags,
        env=env,
        input=stdin,
    )


def judge0_submit(code, api_key, stdin=""):
    payload = json.dumps(
        {"source_code": code, "language_id": 54, "stdin": stdin}
    ).encode()
    req = urllib.request.Request(
        JUDGE0_URL,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "X-RapidAPI-Key": api_key,
            "X-RapidAPI-Host": "judge0-ce.p.rapidapi.com",
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode())
    except (urllib.error.URLError, OSError):
        return None
    status = (result.get("status") or {}).get("description", "Unknown")
    output = result.get("stdout") or ""
    stderr = result.get("stderr") or ""
    compile_out = result.get("compile_output") or ""
    parts = [p for p in (compile_out, output, stderr) if p.strip()]
    text = "\n".join(parts)
    if not text.strip():
        text = "No output."
    if status == "Accepted":
        return jsonify({"ok": True, "output": text, "exit_code": result.get("exit_code", 0)})
    return jsonify(
        {"ok": True, "output": f"[{status}]\n{text}", "exit_code": result.get("exit_code", 1)}
    )


def run_cpp_ce(code, stdin=""):
    payload = json.dumps(
        {
            "source": code,
            "options": {
                "userArguments": "",
                "filters": {"execute": True},
                "executeParameters": {"args": [], "stdin": stdin},
            },
        }
    ).encode()
    req = urllib.request.Request(
        CE_URL,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Accept": "application/json",
            "User-Agent": "CodeFundamentals/1.0",
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read().decode())
    except (urllib.error.URLError, OSError):
        return jsonify(
            {
                "ok": False,
                "error": "Could not reach the online C++ service (Compiler Explorer).",
            }
        )
    if data.get("code") != 0 or data.get("timedOut"):
        err = "\n".join(item.get("text", "") for item in data.get("stderr", []))
        err = re.sub(r"\x1b\[[0-9;]*m", "", err)
        return jsonify(
            {"ok": False, "compile_error": err or "Compilation failed."}
        )
    exec_result = data.get("execResult") or {}
    out_lines = [item.get("text", "") for item in exec_result.get("stdout", [])]
    err_lines = [item.get("text", "") for item in exec_result.get("stderr", [])]
    text = "\n".join([p for p in out_lines + err_lines if p.strip()])
    if exec_result.get("timedOut"):
        text = "Error: execution timed out.\n" + text
    if not text.strip():
        text = "No output."
    return jsonify(
        {"ok": True, "output": text, "exit_code": exec_result.get("code") or 0}
    )


def run_cpp_online(code, stdin=""):
    api_key = os.environ.get("CODE_API_KEY", "")
    if api_key:
        response = judge0_submit(code, api_key, stdin)
        if response is not None:
            return response
    return run_cpp_ce(code, stdin)


def run_java_ce(code, stdin=""):
    ce_code = re.sub(r"\bpublic\s+class\s+Main\b", "class Main", code, count=1)
    payload = json.dumps(
        {
            "source": ce_code,
            "options": {
                "userArguments": "",
                "filters": {"execute": True},
                "executeParameters": {"args": [], "stdin": stdin},
            },
        }
    ).encode()
    req = urllib.request.Request(
        JAVA_CE_URL,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Accept": "application/json",
            "User-Agent": "CodeFundamentals/1.0",
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=90) as resp:
            data = json.loads(resp.read().decode())
    except (urllib.error.URLError, OSError):
        return jsonify(
            {
                "ok": False,
                "error": "Could not reach the online Java service (Compiler Explorer).",
            }
        )
    if data.get("code") != 0 or data.get("timedOut"):
        err = "\n".join(item.get("text", "") for item in data.get("stderr", []))
        err = re.sub(r"\x1b\[[0-9;]*m", "", err)
        return jsonify(
            {"ok": False, "compile_error": err or "Compilation failed."}
        )
    exec_result = data.get("execResult") or {}
    out_lines = [item.get("text", "") for item in exec_result.get("stdout", [])]
    err_lines = [item.get("text", "") for item in exec_result.get("stderr", [])]
    text = "\n".join([p for p in out_lines + err_lines if p.strip()])
    if exec_result.get("timedOut"):
        text = "Error: execution timed out.\n" + text
    if not text.strip():
        text = "No output."
    return jsonify(
        {"ok": True, "output": text, "exit_code": exec_result.get("code") or 0}
    )


def compile_and_run_local(workdir, compile_cmd, stdin=""):
    try:
        compiled = run_locally(compile_cmd, workdir, 30)
    except subprocess.TimeoutExpired:
        return None
    if compiled.returncode != 0:
        return compiled.stderr
    try:
        result = run_locally(
            ["main.exe"], workdir, 10, extra_env=MINGW_BIN, stdin=stdin
        )
    except subprocess.TimeoutExpired:
        return jsonify(
            {
                "ok": True,
                "output": "Error: execution timed out (limit of 10 seconds).",
                "exit_code": -1,
            }
        )
    except (OSError, FileNotFoundError):
        return "BLOCKED"
    output = result.stdout
    if result.stderr.strip():
        output += ("\n" if output.strip() else "") + result.stderr
    return jsonify(
        {"ok": True, "output": output, "exit_code": result.returncode}
    )


def run_cpp(code, stdin=""):
    workdir = tempfile.mkdtemp(prefix="cf_run_")
    try:
        with open(os.path.join(workdir, "main.cpp"), "w", encoding="utf-8") as f:
            f.write(code)

        compile_errors = []
        toolchain_broken = False

        gpp = find_tool("g++")
        if gpp:
            outcome = compile_and_run_local(workdir, [gpp, "main.cpp", "-o", "main.exe"], stdin)
            if isinstance(outcome, dict):
                return outcome
            if outcome == "BLOCKED":
                toolchain_broken = True
            elif outcome is not None:
                if "cannot execute" in outcome or "blocked" in outcome.lower():
                    toolchain_broken = True
                compile_errors.append(outcome)

        clang = find_tool("clang++")
        if clang:
            outcome = compile_and_run_local(
                workdir,
                [
                    clang,
                    "--target=x86_64-w64-windows-gnu",
                    "--sysroot=" + os.path.dirname(MINGW_BIN),
                    "-std=c++17",
                    "main.cpp",
                    "-o",
                    "main.exe",
                ],
                stdin,
            )
            if isinstance(outcome, dict):
                return outcome
            if outcome == "BLOCKED":
                toolchain_broken = True
            elif outcome is not None:
                compile_errors.append(outcome)

        online = run_cpp_online(code, stdin)
        if online is not None:
            return online
        if toolchain_broken:
            return jsonify(
                {
                    "ok": False,
                    "error": (
                        "Local C++ execution is blocked by Windows Smart App Control "
                        "and the online service could not be reached."
                    ),
                }
            ), 501
        message = compile_errors[-1][:4000] if compile_errors else "Compilation failed."
        return jsonify({"ok": False, "compile_error": message})
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def run_ce_compile(ce_url, service_name, code, stdin=""):
    payload = json.dumps(
        {
            "source": code,
            "options": {
                "userArguments": "",
                "filters": {"execute": True},
                "executeParameters": {"args": [], "stdin": stdin},
            },
        }
    ).encode()
    req = urllib.request.Request(
        ce_url,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Accept": "application/json",
            "User-Agent": "CodeFundamentals/1.0",
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read().decode())
    except (urllib.error.URLError, OSError):
        return jsonify(
            {
                "ok": False,
                "error": f"Could not reach the online {service_name} service (Compiler Explorer).",
            }
        )
    if data.get("code") != 0 or data.get("timedOut"):
        err = "\n".join(item.get("text", "") for item in data.get("stderr", []))
        err = re.sub(r"\x1b\[[0-9;]*m", "", err)
        return jsonify(
            {"ok": False, "compile_error": err or "Compilation failed."}
        )
    exec_result = data.get("execResult") or {}
    out_lines = [item.get("text", "") for item in exec_result.get("stdout", [])]
    err_lines = [item.get("text", "") for item in exec_result.get("stderr", [])]
    text = "\n".join([p for p in out_lines + err_lines if p.strip()])
    if exec_result.get("timedOut"):
        text = "Error: execution timed out.\n" + text
    if not text.strip():
        text = "No output."
    return jsonify(
        {"ok": True, "output": text, "exit_code": exec_result.get("code") or 0}
    )


def run_c(code, stdin=""):
    workdir = tempfile.mkdtemp(prefix="cf_run_")
    try:
        with open(os.path.join(workdir, "main.c"), "w", encoding="utf-8") as f:
            f.write(code)

        compile_errors = []
        toolchain_broken = False

        gcc = find_tool("gcc")
        if gcc:
            outcome = compile_and_run_local(workdir, [gcc, "main.c", "-o", "main.exe"], stdin)
            if isinstance(outcome, dict):
                return outcome
            if outcome == "BLOCKED":
                toolchain_broken = True
            elif outcome is not None:
                if "cannot execute" in outcome or "blocked" in outcome.lower():
                    toolchain_broken = True
                compile_errors.append(outcome)

        return run_ce_compile(C_CE_URL, "C", code, stdin)
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


RUN_SESSIONS = {}
MAX_RUN_SECONDS = 120


def _start_local_process(lang, workdir):
    env = dict(os.environ)
    flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
    if lang == "python":
        env["PYTHONUNBUFFERED"] = "1"
        cmd = [sys.executable, "main.py"]
    elif lang == "java":
        cmd = [find_tool("java") or "java", "Main"]
    else:
        env["PATH"] = MINGW_BIN + os.pathsep + env.get("PATH", "")
        cmd = ["main.exe"]
    return subprocess.Popen(
        cmd,
        stdin=subprocess.PIPE,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        cwd=workdir,
        creationflags=flags,
        env=env,
    )


def _pump_pipe(pipe, tag, session):
    try:
        while True:
            ch = pipe.read(1)
            if ch == "":
                break
            session["queue"].put({"type": tag, "data": ch})
    finally:
        session["queue"].put({"type": "done"})


@app.route("/run/start", methods=["POST"])
@login_required
def run_start():
    data = request.get_json(silent=True) or {}
    lang = data.get("language")
    code = data.get("code", "")
    stdin = data.get("stdin", "") or ""
    if lang not in RUN_CONFIGS or not code.strip():
        return jsonify({"ok": False, "error": "Invalid request."}), 400

    if lang == "cpp":
        stripped = re.sub(r"/\*.*?\*/", "", code, flags=re.S)
        stripped = re.sub(r"//[^\n]*", "", stripped)
        if re.search(r"\bcin\b", stripped) and not stdin.strip():
            return jsonify({
                "ok": False,
                "error": "This C++ program reads input with cin. "
                         "Type each value and press Enter."
            }), 400
        result = run_cpp(code, stdin)
        payload = result.get_json()
        events = []
        exit_code = 0
        if payload.get("ok"):
            output = payload.get("output", "")
            if output:
                events.append({"type": "out", "data": output})
            exit_code = payload.get("exit_code", 0)
        else:
            msg = payload.get("compile_error") or payload.get("error") or "Compilation failed."
            events.append({"type": "error", "data": msg})
            exit_code = -1
        return jsonify({
            "ok": True,
            "events": events,
            "finished": True,
            "exit_code": exit_code,
        })

    if lang == "c":
        stripped = re.sub(r"/\*.*?\*/", "", code, flags=re.S)
        stripped = re.sub(r"//[^\n]*", "", stripped)
        if re.search(r"\bscanf\b|\bgets\b|\bfgets\b", stripped) and not stdin.strip():
            return jsonify({
                "ok": False,
                "error": "This C program reads input. "
                         "Type each value and press Enter."
            }), 400
        result = run_c(code, stdin)
        payload = result.get_json()
        events = []
        exit_code = 0
        if payload.get("ok"):
            output = payload.get("output", "")
            if output:
                events.append({"type": "out", "data": output})
            exit_code = payload.get("exit_code", 0)
        else:
            msg = payload.get("compile_error") or payload.get("error") or "Compilation failed."
            events.append({"type": "error", "data": msg})
            exit_code = -1
        return jsonify({
            "ok": True,
            "events": events,
            "finished": True,
            "exit_code": exit_code,
        })

    workdir = tempfile.mkdtemp(prefix="cf_run_")
    source = RUN_CONFIGS[lang]["source"]
    try:
        with open(os.path.join(workdir, source), "w", encoding="utf-8") as f:
            f.write(code)
        events = []
        exit_code = 0

        if lang == "java":
            javac = find_tool("javac")
            if not javac:
                result = run_java_ce(code, stdin)
                payload = result.get_json()
                if not payload.get("ok"):
                    return jsonify(
                        {"ok": False, "compile_error": payload.get("compile_error") or payload.get("error")}
                    ), 502
                events = []
                output = payload.get("output", "")
                if output:
                    events.append({"type": "out", "data": output})
                return jsonify({
                    "ok": True,
                    "events": events,
                    "finished": True,
                    "exit_code": payload.get("exit_code", 0),
                })
            compiled = run_locally([javac, "Main.java"], workdir, 20)
            if compiled.returncode != 0:
                return jsonify({"ok": False, "compile_error": compiled.stderr[:4000]})
            run_cmd = [find_tool("java") or "java", "Main"]
        elif lang == "typescript":
            node = shutil.which("node")
            if not node:
                return jsonify(
                    {
                        "ok": False,
                        "compile_error": "TypeScript needs Node.js on the server, but it was not found.",
                    }
                ), 502
            run_cmd = [node, "--experimental-strip-types", "--no-warnings", "main.ts"]
        else:
            run_cmd = [sys.executable, "main.py"]

        try:
            result = run_locally(run_cmd, workdir, MAX_RUN_SECONDS, stdin=stdin)
        except subprocess.TimeoutExpired:
            events.append({
                "type": "error",
                "data": "Program was stopped (time limit of %d seconds)." % MAX_RUN_SECONDS,
            })
            return jsonify({
                "ok": True, "events": events, "finished": True, "exit_code": -1,
            })

        output = result.stdout
        if result.stderr.strip():
            output += ("\n" if output.strip() else "") + result.stderr
        if output.strip():
            events.append({"type": "out", "data": output})
        exit_code = result.returncode
        return jsonify({
            "ok": True, "events": events, "finished": True, "exit_code": exit_code,
        })
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


@app.route("/run/poll/<session_id>", methods=["POST"])
@login_required
def run_poll(session_id):
    session = RUN_SESSIONS.get(session_id)
    if not session:
        return jsonify({"ok": False, "error": "Session not found."}), 404
    proc = session.get("proc")
    events = []
    exit_code = None

    started = session.get("started")
    if started and time.time() - started > MAX_RUN_SECONDS:
        if proc and proc.poll() is None:
            proc.kill()
        session["queue"].put(
            {
                "type": "error",
                "data": "Program was stopped (time limit of %d seconds)."
                % MAX_RUN_SECONDS,
            }
        )
        session["queue"].put({"type": "exit", "data": -1})

    done_count = session.get("done_count", 0)
    text_buf = ""
    while True:
        try:
            item = session["queue"].get_nowait()
        except queue.Empty:
            break
        if item["type"] == "exit":
            exit_code = item["data"]
        elif item["type"] == "done":
            done_count += 1
        elif item["type"] == "error":
            if text_buf:
                events.append({"type": "out", "data": text_buf.rstrip("\r\n")})
                text_buf = ""
            events.append(item)
        else:
            text_buf += item["data"]
    if text_buf:
        events.append({"type": "out", "data": text_buf.rstrip("\r\n")})
    session["done_count"] = done_count

    if exit_code is None and proc is not None:
        if done_count >= 2 and proc.poll() is not None:
            exit_code = proc.returncode
        elif proc.poll() is not None:
            if session.get("drain_until") is None:
                session["drain_until"] = time.time() + 1.5
            elif time.time() > session["drain_until"]:
                exit_code = proc.returncode
        else:
            session.pop("drain_until", None)

    if exit_code is not None:
        if proc:
            try:
                proc.kill()
            except Exception:
                pass
        shutil.rmtree(session.get("workdir", ""), ignore_errors=True)
        RUN_SESSIONS.pop(session_id, None)
        return jsonify(
            {"ok": True, "events": events, "finished": True, "exit_code": exit_code}
        )
    return jsonify({"ok": True, "events": events, "finished": False})


@app.route("/run/input/<session_id>", methods=["POST"])
@login_required
def run_input(session_id):
    session = RUN_SESSIONS.get(session_id)
    if not session or not session.get("proc"):
        return jsonify({"ok": False, "error": "Session is not running."}), 404
    data = request.get_json(silent=True) or {}
    text = data.get("input", "")
    try:
        session["proc"].stdin.write(text + "\n")
        session["proc"].stdin.flush()
    except (BrokenPipeError, OSError, ValueError):
        return jsonify(
            {"ok": False, "error": "The program is no longer accepting input."}
        ), 400
    return jsonify({"ok": True})


@app.route("/run/stop/<session_id>", methods=["POST"])
@login_required
def run_stop(session_id):
    session = RUN_SESSIONS.get(session_id)
    if not session:
        return jsonify({"ok": True})
    proc = session.get("proc")
    if proc and proc.poll() is None:
        proc.kill()
    return jsonify({"ok": True})


def format_sql_table(cols, rows):
    if not rows:
        return f"-- {len(cols)} column(s), 0 rows returned"
    widths = [len(c) for c in cols]
    for row in rows:
        for i, value in enumerate(row):
            widths[i] = max(widths[i], len(str(value)))
    sep = "+" + "+".join("-" * (w + 2) for w in widths) + "+"

    def line(row):
        cells = []
        for i, value in enumerate(row):
            cells.append(" " + str(value).ljust(widths[i]) + " ")
        return "|" + "|".join(cells) + "|"

    lines = [sep, line(cols), sep]
    for row in rows:
        lines.append(line(row))
    lines.append(sep)
    label = f"({len(rows)} row returned)" if len(rows) == 1 else f"({len(rows)} rows returned)"
    lines.append(label)
    return "\n".join(lines)


@app.route("/run/sql", methods=["POST"])
@login_required
def run_sql():
    data = request.get_json(silent=True) or {}
    code = (data.get("code") or "").strip()
    if not code:
        return jsonify({"ok": False, "error": "No SQL provided."}), 400
    if len(code) > 50000:
        return jsonify({"ok": False, "error": "SQL is too long (50,000 characters max)."}), 400
    conn = sqlite3.connect(":memory:")
    output = []
    try:
        for statement in code.split(";"):
            statement = statement.strip()
            if not statement:
                continue
            cur = conn.execute(statement)
            if cur.description:
                output.append(format_sql_table([d[0] for d in cur.description], cur.fetchall()))
            elif cur.rowcount != -1:
                output.append(f"-- {cur.rowcount} row(s) affected")
            else:
                output.append("-- Statement executed")
        return jsonify({"ok": True, "output": output})
    except sqlite3.Error as e:
        return jsonify({"ok": False, "error": str(e)}), 400
    finally:
        conn.close()


PLAYGROUND_SCHEMA = [
    "CREATE TABLE users (id INTEGER PRIMARY KEY, username TEXT, password TEXT, role TEXT)",
    "INSERT INTO users (id, username, password, role) VALUES "
    "(1, 'admin', 'admin123', 'admin'), "
    "(2, 'alice', 'wonderland', 'staff'), "
    "(3, 'bob', 'builder', 'staff'), "
    "(4, 'carol', 'carol123', 'guest')",
    "CREATE TABLE products (id INTEGER PRIMARY KEY, name TEXT, price INTEGER)",
    "INSERT INTO products (id, name, price) VALUES "
    "(1, 'Laptop', 45000), "
    "(2, 'Mouse', 500), "
    "(3, 'Keyboard', 1200), "
    "(4, 'Monitor', 8000)",
]


@app.route("/sql-playground")
@login_required
def sql_playground():
    return render_template("sql_playground.html")


@app.route("/run/playground", methods=["POST"])
@login_required
def run_playground():
    data = request.get_json(silent=True) or {}
    code = (data.get("code") or "").strip()
    if not code:
        return jsonify({"ok": False, "error": "No SQL provided."}), 400
    if len(code) > 50000:
        return jsonify({"ok": False, "error": "SQL is too long (50,000 characters max)."}), 400
    conn = sqlite3.connect(":memory:")
    output = []
    try:
        for statement in PLAYGROUND_SCHEMA:
            conn.execute(statement)
        for statement in code.split(";"):
            statement = statement.strip()
            if not statement:
                continue
            cur = conn.execute(statement)
            if cur.description:
                output.append(format_sql_table([d[0] for d in cur.description], cur.fetchall()))
            elif cur.rowcount != -1:
                output.append(f"-- {cur.rowcount} row(s) affected")
            else:
                output.append("-- Statement executed")
        return jsonify({"ok": True, "output": output})
    except sqlite3.Error as e:
        return jsonify({"ok": False, "error": str(e)}), 400
    finally:
        conn.close()


# ==================== DATABASE FIX FUNCTION ====================
def fix_database_schema():
    """Fix missing columns in the database schema."""
    from sqlalchemy import inspect as _inspect
    
    with app.app_context():
        inspector = _inspect(db.engine)
        is_sqlite = db.engine.dialect.name == "sqlite"
        
        try:
            # Check if section_quiz_score table exists
            tables = inspector.get_table_names()
            if "section_quiz_score" not in tables:
                print("[FIX] section_quiz_score table does not exist, creating...")
                db.create_all()
                return
            
            # Get existing columns
            columns = [col['name'] for col in inspector.get_columns('section_quiz_score')]
            
            # Add topic_ratings column if missing
            if 'topic_ratings' not in columns:
                print("[FIX] Adding topic_ratings column to section_quiz_score...")
                with db.engine.begin() as conn:
                    if is_sqlite:
                        conn.execute(sa_text("ALTER TABLE section_quiz_score ADD COLUMN topic_ratings TEXT"))
                    else:
                        conn.execute(sa_text("ALTER TABLE section_quiz_score ADD COLUMN topic_ratings TEXT"))
                print("[FIX] topic_ratings column added successfully!")
            else:
                print("[FIX] topic_ratings column already exists.")
                
        except Exception as e:
            print(f"[FIX ERROR] Failed to fix database schema: {e}")


# ==================== FIX ROUTE ====================
@app.route("/fix-db")
@login_required
@admin_required
def fix_db_route():
    """Admin route to fix database schema issues."""
    try:
        fix_database_schema()
        flash("Database schema fixed successfully! topic_ratings column has been added.", "success")
    except Exception as e:
        flash(f"Error fixing database: {e}", "danger")
    return redirect(url_for("admin_dashboard"))


# ==================== STARTUP ====================
with app.app_context():
    # Fix database schema first
    fix_database_schema()
    
    db.create_all()
    os.makedirs(AVATAR_DIR, exist_ok=True)
    is_sqlite = db.engine.dialect.name == "sqlite"
    
    # Check if user table has required columns, drop and recreate if corrupted
    from sqlalchemy import inspect as _inspect
    try:
        _user_cols = [c["name"] for c in _inspect(db.engine).get_columns("user")]
        if "username" not in _user_cols:
            # Table is corrupted, drop and recreate all tables
            with db.engine.begin() as conn:
                if is_sqlite:
                    conn.execute(sa_text("DROP TABLE IF EXISTS lesson_progress"))
                    conn.execute(sa_text("DROP TABLE IF EXISTS quiz_score"))
                    conn.execute(sa_text("DROP TABLE IF EXISTS section_quiz_score"))
                    conn.execute(sa_text("DROP TABLE IF EXISTS lab_progress"))
                    conn.execute(sa_text("DROP TABLE IF EXISTS activity_log"))
                    conn.execute(sa_text("DROP TABLE IF EXISTS feedback"))
                    conn.execute(sa_text("DROP TABLE IF EXISTS user"))
                else:
                    conn.execute(sa_text('DROP TABLE IF EXISTS lesson_progress CASCADE'))
                    conn.execute(sa_text('DROP TABLE IF EXISTS quiz_score CASCADE'))
                    conn.execute(sa_text('DROP TABLE IF EXISTS section_quiz_score CASCADE'))
                    conn.execute(sa_text('DROP TABLE IF EXISTS lab_progress CASCADE'))
                    conn.execute(sa_text('DROP TABLE IF EXISTS activity_log CASCADE'))
                    conn.execute(sa_text('DROP TABLE IF EXISTS feedback CASCADE'))
                    conn.execute(sa_text('DROP TABLE IF EXISTS "user" CASCADE'))
            db.create_all()
    except Exception:
        pass
    
    if is_sqlite:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE user ADD COLUMN avatar VARCHAR(200)"))
                break
            except Exception:
                time.sleep(0.5)
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE user ADD COLUMN lab_completed INTEGER DEFAULT 0"))
                break
            except Exception:
                time.sleep(0.5)
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE section_quiz_score ADD COLUMN last_answers TEXT"))
                break
            except Exception:
                time.sleep(0.5)
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE section_quiz_score ADD COLUMN topic_ratings TEXT"))
                break
            except Exception:
                time.sleep(0.5)
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE user ADD COLUMN coins INTEGER DEFAULT 0"))
                break
            except Exception:
                time.sleep(0.5)
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE user ADD COLUMN is_banned BOOLEAN DEFAULT 0"))
                break
            except Exception:
                time.sleep(0.5)
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text("ALTER TABLE user ADD COLUMN ban_reason VARCHAR(200)"))
                break
            except Exception:
                time.sleep(0.5)
    from sqlalchemy import inspect as _inspect

    try:
        _cols = [c["name"] for c in _inspect(db.engine).get_columns("user")]
    except Exception:
        _cols = []
    if "avatar" in _cols and not is_sqlite:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text('ALTER TABLE "user" ALTER COLUMN avatar TYPE TEXT'))
                break
            except Exception:
                time.sleep(0.5)
    if "is_admin" not in _cols:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    if is_sqlite:
                        conn.execute(sa_text("ALTER TABLE user ADD COLUMN is_admin BOOLEAN DEFAULT 0"))
                    else:
                        conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN is_admin BOOLEAN DEFAULT FALSE'))
                break
            except Exception:
                time.sleep(0.5)
    if "last_seen" not in _cols:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    if is_sqlite:
                        conn.execute(sa_text("ALTER TABLE user ADD COLUMN last_seen DATETIME"))
                    else:
                        conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN last_seen TIMESTAMP'))
                break
            except Exception:
                time.sleep(0.5)
    if "coins" not in _cols and not is_sqlite:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN coins INTEGER DEFAULT 0'))
                break
            except Exception:
                time.sleep(0.5)
    if "is_banned" not in _cols:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN is_banned BOOLEAN DEFAULT FALSE'))
                break
            except Exception:
                time.sleep(0.5)
    if "ban_reason" not in _cols:
        for _ in range(10):
            try:
                with db.engine.begin() as conn:
                    conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN ban_reason VARCHAR(200)'))
                break
            except Exception:
                time.sleep(0.5)
    admin_username = os.environ.get("ADMIN_USERNAME", "").strip()
    try:
        if admin_username:
            admin_user = User.query.filter_by(username=admin_username).first()
            if admin_user and not admin_user.is_admin:
                admin_user.is_admin = True
                db.session.commit()
        admin_acct = User.query.filter_by(username="admin").first()
        if admin_acct is None:
            db.session.add(
                User(
                    username="admin",
                    email="admin@codefundamentals.local",
                    password_hash=generate_password_hash("admin123"),
                    is_admin=True,
                )
            )
            db.session.commit()
        elif not admin_acct.is_admin:
            admin_acct.is_admin = True
            db.session.commit()
    except Exception:
        pass
    try:
        reset_marker = db.session.execute(
            sa_text("SELECT COUNT(*) FROM activity_log WHERE action = :a"),
            {"a": "SYSTEM: progress reset v1.48"},
        ).scalar()
        if not reset_marker:
            db.session.query(LessonProgress).filter(
                LessonProgress.lesson_id >= 5
            ).delete(synchronize_session=False)
            db.session.query(QuizScore).filter(
                QuizScore.lesson_id >= 5
            ).delete(synchronize_session=False)
            anchor = User.query.order_by(User.id).first()
            if anchor is not None:
                db.session.add(
                    ActivityLog(
                        user_id=anchor.id,
                        action="SYSTEM: progress reset v1.48",
                        details="Reset shifted lesson/quiz progress from legacy migration",
                    )
                )
            db.session.commit()
    except Exception:
        db.session.rollback()
    try:
        legacy_users = db.session.execute(
            sa_text(
                "SELECT id FROM user WHERE lab_completed > 0 "
                "AND NOT EXISTS (SELECT 1 FROM lab_progress lp WHERE lp.user_id = user.id)"
            )
        ).fetchall()
        for (uid,) in legacy_users:
            db.session.add(LabProgress(user_id=uid, lesson_id=19))
        if legacy_users:
            db.session.commit()
    except Exception:
        pass


if __name__ == "__main__":
    app.run(host="0.0.0.0", debug=True, threaded=True, port=int(os.environ.get("PORT", 5001)))